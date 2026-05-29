import streamlit as st
import pandas as pd
from datetime import datetime, timedelta, date
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
from collections import Counter
import matplotlib.pyplot as plt
from openpyxl.drawing.image import Image as XLImage
from openpyxl import Workbook
from openpyxl.chart import PieChart, Reference
from openpyxl.chart.series import DataPoint
from openpyxl.styles import PatternFill
from openpyxl.chart.label import DataLabelList
import numpy as np
import unicodedata
import pytz
import re

# --- CONFIGURACIÓN DE ZONA HORARIA ---
ZONA_HORARIA_ARGENTINA = pytz.timezone('America/Argentina/Buenos_Aires')

def obtener_fecha_actual_argentina():
    """Obtiene la fecha actual en la zona horaria de Argentina"""
    return datetime.now(ZONA_HORARIA_ARGENTINA)

# --- FUNCIÓN PARA GENERAR EXCEL (MOVIDA AL INICIO) ---
def generar_excel_desde_df(df, nombre_hoja="Datos"):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        df.to_excel(writer, sheet_name=nombre_hoja, index=False)
    output.seek(0)
    return output

# --- SISTEMA MEJORADO DE FERIADOS Y PUENTES ---

def es_dia_festivo(fecha=None):
    """Verifica si la fecha es un día festivo configurado"""
    if fecha is None:
        fecha = date.today()
    if isinstance(fecha, datetime):
        fecha = fecha.date()
    festivos = [
        (1, 1),   # Año Nuevo
        (2, 16),  # Carnaval
        (2, 17),  # Carnaval
        (3, 23),  # Feriado Puente
        (3, 24),  # Día Nacional de la Memoria
        (4, 2),   # Día del Veterano
        (4, 3),   # Viernes Santo
        (5, 1),   # Día del Trabajo
        (5, 25),  # Día de la Revolución de Mayo
        (6, 17),  # Paso a la Inmortalidad del Gral. Martín Güemes
        (6, 20),  # Día de la Bandera
        (7, 9),   # Día de la Independencia
        (10, 12), # Día de la Raza
        (11, 21), # Día de la Soberanía Nacional
        (11, 24), # Día de la Virgen
        (12, 8),  # Inmaculada Concepción
        (12, 24), # Pre Navidad
        (12, 25), # Navidad
        (12, 31), # Fin de Año
    ]
    return (fecha.month, fecha.day) in festivos

def es_feriado_puente(fecha):
    """Detecta feriados puente con mensaje descriptivo."""
    if isinstance(fecha, datetime):
        fecha = fecha.date()
    if fecha.weekday() == 4:  # Viernes
        sabado = fecha + timedelta(days=1)
        domingo = fecha + timedelta(days=2)
        if es_dia_festivo(sabado) and es_dia_festivo(domingo):
            return True, f"Viernes puente (festivo {sabado.strftime('%d/%m')} y {domingo.strftime('%d/%m')})"
        elif es_dia_festivo(sabado):
            return True, f"Viernes puente (festivo {sabado.strftime('%d/%m')})"
        elif es_dia_festivo(domingo):
            return True, f"Viernes puente (festivo {domingo.strftime('%d/%m')})"
    return False, ""

def es_dia_laborable(fecha):
    """Determina si una fecha es laborable (no fin de semana, no feriado, no puente)"""
    if isinstance(fecha, datetime):
        fecha = fecha.date()
    if fecha.weekday() >= 5:
        return False
    if es_dia_festivo(fecha):
        return False
    if es_feriado_puente(fecha)[0]:
        return False
    return True

def es_dia_habil(fecha):
    """Determina si un día es hábil (versión mejorada con puentes)"""
    return es_dia_laborable(fecha)

def es_feriado(fecha):
    """Determina si un día es feriado (versión mejorada)"""
    return es_dia_festivo(fecha) or es_feriado_puente(fecha)[0]

# --- DICCIONARIO DE SEMANAS REALES (CALENDARIO) ---
def obtener_semana_calendario(fecha):
    """Calcula la semana del año según calendario ISO 8601"""
    if pd.isna(fecha):
        return None
    try:
        if isinstance(fecha, str):
            fecha = pd.to_datetime(fecha, errors='coerce')
        semana = fecha.isocalendar()[1]
        return semana
    except:
        return None

# --- FUNCIÓN CALCULAR DÍAS HÁBILES ---
def calcular_dias_habiles(fecha_inicio, fecha_fin):
    """Calcula días hábiles entre dos fechas considerando feriados y puentes"""
    if pd.isna(fecha_inicio) or pd.isna(fecha_fin):
        return None
    if hasattr(fecha_inicio, 'tz') and fecha_inicio.tz is not None:
        fecha_inicio = fecha_inicio.replace(tzinfo=None)
    if hasattr(fecha_fin, 'tz') and fecha_fin.tz is not None:
        fecha_fin = fecha_fin.replace(tzinfo=None)
    fecha_inicio = fecha_inicio.date() if hasattr(fecha_inicio, 'date') else fecha_inicio
    fecha_fin = fecha_fin.date() if hasattr(fecha_fin, 'date') else fecha_fin
    if fecha_inicio > fecha_fin:
        return 0
    dias = 0
    current = fecha_inicio + timedelta(days=1)
    while current <= fecha_fin:
        if es_dia_habil(current):
            dias += 1
        current += timedelta(days=1)
    return dias

def calcular_conteo_diario_amba(df):
    """Genera conteo diario de AMBA y detalle de pendientes, con agrupación correcta."""
    # --- Filtrado ---
    mask_amba = df['Categoria'].isin(['AMBA cercano', 'AMBA interior'])
    mask_excluir = (
        df['Cumplimiento'].isin(['Cancelada', 'Excluido - Logística Inversa POS']) |
        df['Cumplimiento'].str.contains('Devuelto', na=False)
    )
    df_amba = df[mask_amba & ~mask_excluir].copy()

    # --- Columnas de trabajo ---
    df_amba['es_entregado'] = df_amba['Estado'].str.contains('Entregada', na=False)
    df_amba['sla_dias'] = df_amba['Categoria'].map({'AMBA cercano': 2, 'AMBA interior': 5})
    df_amba['lt_num'] = pd.to_numeric(df_amba['Lead Time'], errors='coerce')
    df_amba['a_tiempo'] = df_amba['es_entregado'] & (df_amba['lt_num'] <= df_amba['sla_dias'])

    # --- Fecha de agrupación (siguiente día hábil si no está entregado) ---
    df_amba['fecha_ref'] = pd.to_datetime(df_amba['Fecha'], errors='coerce')

    def siguiente_dia_habil(fecha):
        if pd.isnull(fecha):
            return pd.NaT
        d = fecha + timedelta(days=1)
        while not es_dia_habil(d):
            d += timedelta(days=1)
        return d

    df_amba['fecha_agrup'] = np.where(
        df_amba['es_entregado'],
        df_amba['fecha_ref'],
        df_amba['fecha_ref'].apply(siguiente_dia_habil)
    )

    # --- CLAVE: normalizar fecha a solo día, sin hora ni timezone ---
    df_amba['fecha_agrup_d'] = pd.to_datetime(df_amba['fecha_agrup']).dt.normalize().dt.tz_localize(None)

    # --- Variables numéricas para agregación ---
    df_amba['entregado_int'] = df_amba['es_entregado'].astype(int)
    df_amba['tiempo_int'] = df_amba['a_tiempo'].astype(int)
    df_amba['fuera_int'] = (df_amba['es_entregado'] & ~df_amba['a_tiempo']).astype(int)
    df_amba['pendiente_int'] = (~df_amba['es_entregado']).astype(int)

    # --- Agrupación diaria por categoría ---
    conteos = df_amba.groupby(['fecha_agrup_d', 'Categoria'], as_index=False).agg(
        Entregados_En_Tiempo=('tiempo_int', 'sum'),
        Entregados_Fuera=('fuera_int', 'sum'),
        Total_Entregados=('entregado_int', 'sum'),
        Total_Pendientes=('pendiente_int', 'sum')
    )

    conteos['Total_Gestionable'] = conteos['Total_Entregados'] + conteos['Total_Pendientes']
    conteos['% Entregas en Tiempo'] = (conteos['Entregados_En_Tiempo'] / conteos['Total_Gestionable']).fillna(0)
    conteos['% Entregas Totales'] = (conteos['Total_Entregados'] / conteos['Total_Gestionable']).fillna(0)

    # Renombrar y ordenar
    conteos = conteos.rename(columns={'fecha_agrup_d': 'Fecha', 'Categoria': 'Categoría'})
    conteos = conteos.sort_values(['Fecha', 'Categoría']).reset_index(drop=True)

    # Formatear fecha para visualización
    conteos['Fecha'] = conteos['Fecha'].dt.strftime('%d/%m/%Y')

    # --- Detalle de pendientes ---
    df_pend = df_amba[~df_amba['es_entregado']].copy()
    df_pend['Demora'] = df_pend['lt_num'] - df_pend['sla_dias']
    df_pend['Demora'] = df_pend['Demora'].apply(lambda x: f"+{int(x)} días" if pd.notna(x) else "-")
    df_pend['Observación'] = df_pend['Cumplimiento']

    detalle = df_pend[[
        'fecha_agrup_d', 'Guia', 'Destinatario', 'Localidad destino',
        'Categoria', 'Días Prometidos', 'Lead Time', 'Visitas', 'Demora', 'Observación'
    ]].copy()

    detalle = detalle.rename(columns={
        'fecha_agrup_d': 'Fecha',
        'Categoria': 'Categoría',
        'Localidad destino': 'Localidad',
        'Días Prometidos': 'Días Prom.',
        'Lead Time': 'Lead Time'
    })
    detalle['Fecha'] = detalle['Fecha'].dt.strftime('%d/%m/%Y')
    detalle = detalle.sort_values(['Fecha', 'Categoría']).reset_index(drop=True)

    return conteos, detalle


# --- LISTA DE LOCALIDADES AMBA ---
amba_localidades = [
    "CIUDAD AUTONOMA BUENOS AIRES", "AVELLANEDA", "LANUS", "LOMAS DE ZAMORA",
    "LA MATANZA", "MORON", "SAN MARTIN", "VICENTE LOPEZ", "SAN ISIDRO",
    "TRES DE FEBRERO", "MORENO", "HURLINGHAM", "ITUZAINGO", "BERAZATEGUI",
    "FLORENCIO VARELA", "QUILMES", "ALMIRANTE BROWN", "ESTEBAN ECHEVERRIA",
    "EZEIZA", "SAN FERNANDO", "TIGRE", "SAN MIGUEL", "MALVINAS ARGENTINAS",
    "JOSE C. PAZ", "PILAR", "ESCOBAR", "MERLO", "MARCOS PAZ",
    "GENERAL RODRIGUEZ", "PRESIDENTE PERON", "SAN VICENTE", "BRANDSEN",
    "BERISSO", "ENSENADA", "LA PLATA", "MUNRO", "SAAVEDRA", "FLORES",
    "ALMAGRO", "VILLA URQUIZA", "COLEGIALES", "PALERMO", "RECOLETA",
    "BELGRANO", "NUÑEZ", "CABALLITO", "BOEDO", "SAN TELMO", "CONSTITUCION",
    "RETIRO", "SAN CRISTOBAL", "BALVANERA", "MONTSERRAT",
    "CAÑUELAS, BUENOS AIRES", "ZARATE, BUENOS AIRES",
    "LOMAS DEL MIRADOR , BUENOS AIRES", "MORENO MARIANO, BUENOS AIRES"
]

excepciones_amba = [
    "SAN MARTIN, SANTA FE", "SAN MARTIN, MENDOZA", "SAN MARTIN, SAN JUAN",
    "SAN MARTIN, CORRIENTES", "SAN MARTIN, ENTRE RIOS",
    "VILLA LIB. GENERAL SAN MARTIN", "GENERAL SAN MARTIN",
    "SAN MARTIN DE LOS ANDES", "SAN MARTIN DE LA VEGA",
    "TANDIL, BUENOS AIRES", "MAR DEL PLATA, BUENOS AIRES",
    "BAHIA BLANCA, BUENOS AIRES", "NECOCHEA, BUENOS AIRES",
    "OLAVARRIA, BUENOS AIRES", "AZUL, BUENOS AIRES",
    "SAN MIGUEL DE TUCUMAN", "SAN MIGUEL, TUCUMAN", "TUCUMAN, TUCUMAN"
]

def determinar_zona(localidad_destino):
    """Determina si una localidad pertenece a AMBA o INTERIOR"""
    if pd.isna(localidad_destino) or localidad_destino == "":
        return "INTERIOR"
    localidad = str(localidad_destino).upper().strip()

    def normalizar_texto(texto):
        return ''.join(
            c for c in unicodedata.normalize('NFD', texto)
            if unicodedata.category(c) != 'Mn'
        ).replace("Ñ", "N")

    localidad_normalizada = normalizar_texto(localidad)

    provincias_interior = [
        "TUCUMAN", "CATAMARCA", "LA RIOJA", "SANTIAGO DEL ESTERO", "SALTA",
        "JUJUY", "MENDOZA", "SAN JUAN", "SAN LUIS", "CORDOBA", "SANTA FE",
        "ENTRE RIOS", "CORRIENTES", "MISIONES", "CHACO", "FORMOSA", "NEUQUEN",
        "RIO NEGRO", "CHUBUT", "SANTA CRUZ", "TIERRA DEL FUEGO", "LA PAMPA"
    ]

    for provincia in provincias_interior:
        provincia_normalizada = normalizar_texto(provincia)
        if provincia_normalizada in localidad_normalizada:
            patron_provincia = r'\b' + re.escape(provincia_normalizada) + r'\b'
            if re.search(patron_provincia, localidad_normalizada):
                return "INTERIOR"

    excepciones_amba_actualizadas = [
        "SAN MARTIN, SANTA FE", "SAN MARTIN, MENDOZA", "SAN MARTIN, SAN JUAN",
        "SAN MARTIN, CORRIENTES", "SAN MARTIN, ENTRE RIOS",
        "VILLA LIB. GENERAL SAN MARTIN", "GENERAL SAN MARTIN",
        "SAN MARTIN DE LOS ANDES", "SAN MARTIN DE LA VEGA",
        "TANDIL, BUENOS AIRES", "MAR DEL PLATA, BUENOS AIRES",
        "BAHIA BLANCA, BUENOS AIRES", "NECOCHEA, BUENOS AIRES",
        "OLAVARRIA, BUENOS AIRES", "AZUL, BUENOS AIRES",
        "SAN MIGUEL DE TUCUMAN", "SAN MIGUEL, TUCUMAN", "TUCUMAN, TUCUMAN",
        "YERBA BUENA, TUCUMAN", "TAFI VIEJO, TUCUMAN",
        "LAS TALITAS, TUCUMAN", "BANDA DEL RIO SALI, TUCUMAN"
    ]

    for excepcion in excepciones_amba_actualizadas:
        excepcion_normalizada = normalizar_texto(excepcion)
        if excepcion_normalizada == localidad_normalizada:
            return "INTERIOR"
        if excepcion_normalizada in localidad_normalizada:
            patron_excepcion = r'\b' + re.escape(excepcion_normalizada) + r'\b'
            if re.search(patron_excepcion, localidad_normalizada):
                return "INTERIOR"

    for localidad_amba in amba_localidades:
        amba_normalizada = normalizar_texto(localidad_amba)
        if amba_normalizada == localidad_normalizada:
            return "AMBA"
        if localidad_normalizada.startswith(amba_normalizada):
            resto = localidad_normalizada[len(amba_normalizada):].strip()
            if len(resto) == 0:
                return "AMBA"
            if resto.startswith(",") or resto.startswith(" "):
                resto_limpio = resto.lstrip(", ").upper()
                if "BUENOS AIRES" in resto_limpio or len(resto_limpio) == 0:
                    return "AMBA"
                if any(provincia in resto_limpio for provincia in provincias_interior):
                    return "INTERIOR"

    palabras_caba = ["CAPITAL FEDERAL", "C.A.B.A.", "CABA", "CIUDAD AUTONOMA"]
    for palabra in palabras_caba:
        if palabra in localidad:
            return "AMBA"

    return "INTERIOR"

def limpiar_localidad(localidad):
    """Limpia las localidades eliminando partes duplicadas"""
    if pd.isna(localidad):
        return localidad
    loc_str = str(localidad).upper().strip()
    patron_est = r',\s*EST\.[^,]*'
    loc_str = re.sub(patron_est, '', loc_str, flags=re.IGNORECASE)
    patron_est_solo = r',\s*EST[,\s]*'
    loc_str = re.sub(patron_est_solo, ',', loc_str, flags=re.IGNORECASE)
    loc_str = re.sub(r',\s*,', ',', loc_str)
    loc_str = re.sub(r',\s*$', '', loc_str)
    loc_str = re.sub(r'^\s*,\s*', '', loc_str)
    loc_str = re.sub(r'\s+', ' ', loc_str)
    loc_str = loc_str.strip().strip(',')
    if ',' in loc_str:
        partes = [p.strip() for p in loc_str.split(',') if p.strip()]
        if len(partes) == 3:
            nombre_localidad = f"{partes[0]} {partes[1]}"
            provincia = partes[2]
            return f"{nombre_localidad}, {provincia}"
        elif len(partes) == 2:
            return f"{partes[0]}, {partes[1]}"
        elif len(partes) == 1:
            return partes[0]
    return loc_str

def determinar_categoria(localidad_destino):
    """Clasifica una localidad en: AMBA cercano, AMBA interior, Buenos Aires interior, Interior"""
    if pd.isna(localidad_destino) or localidad_destino == "":
        return "Interior"
    loc_str = str(localidad_destino).upper().strip()

    def normalizar(texto):
        return ''.join(c for c in unicodedata.normalize('NFD', texto) if unicodedata.category(c) != 'Mn')

    loc_norm = normalizar(loc_str)

    provincias_interior = [
        "TUCUMAN", "CATAMARCA", "LA RIOJA", "SANTIAGO DEL ESTERO", "SALTA",
        "JUJUY", "MENDOZA", "SAN JUAN", "SAN LUIS", "CORDOBA", "SANTA FE",
        "ENTRE RIOS", "CORRIENTES", "MISIONES", "CHACO", "FORMOSA", "NEUQUEN",
        "RIO NEGRO", "CHUBUT", "SANTA CRUZ", "TIERRA DEL FUEGO", "LA PAMPA"
    ]
    for prov in provincias_interior:
        if re.search(r'\b' + re.escape(normalizar(prov)) + r'\b', loc_norm):
            return "Interior"

    es_buenos_aires = "BUENOS AIRES" in loc_norm

    amba_lejano_lista = [
        "GENERAL RODRIGUEZ", "CITY BELL", "TOLOSA", "CAÑUELAS", "LA PLATA",
        "PILAR", "GENERAL LAS HERAS", "BRANDSEN", "LUJAN", "CAMPANA", "ZARATE"
    ]
    for ciudad in amba_lejano_lista:
        if re.search(r'\b' + re.escape(normalizar(ciudad)) + r'\b', loc_norm):
            if es_buenos_aires or ciudad in ["CITY BELL", "TOLOSA", "LA PLATA"]:
                return "AMBA interior"

    amba_cercano_lista = [
        "CIUDAD AUTONOMA BUENOS AIRES", "CAPITAL FEDERAL", "CABA",
        "AVELLANEDA", "LANUS", "LOMAS DE ZAMORA", "LA MATANZA", "MORON",
        "SAN MARTIN", "VICENTE LOPEZ", "SAN ISIDRO", "TRES DE FEBRERO",
        "MORENO", "HURLINGHAM", "ITUZAINGO", "BERAZATEGUI", "FLORENCIO VARELA",
        "QUILMES", "ALMIRANTE BROWN", "ESTEBAN ECHEVERRIA", "EZEIZA",
        "SAN FERNANDO", "TIGRE", "SAN MIGUEL", "MALVINAS ARGENTINAS",
        "JOSE C. PAZ", "ESCOBAR", "MERLO", "MARCOS PAZ", "PRESIDENTE PERON",
        "SAN VICENTE", "BERISSO", "ENSENADA", "MUNRO", "SAAVEDRA", "FLORES",
        "ALMAGRO", "VILLA URQUIZA", "COLEGIALES", "PALERMO", "RECOLETA",
        "BELGRANO", "NUÑEZ", "CABALLITO", "BOEDO", "SAN TELMO", "CONSTITUCION",
        "RETIRO", "SAN CRISTOBAL", "BALVANERA", "MONTSERRAT",
        "LOMAS DEL MIRADOR", "MORENO MARIANO"
    ]
    for ciudad in amba_cercano_lista:
        if re.search(r'\b' + re.escape(normalizar(ciudad)) + r'\b', loc_norm):
            return "AMBA cercano"

    if es_buenos_aires:
        return "Buenos Aires interior"

    return "Interior"

# --- CONSTANTE: AGENCIA ORIGEN CORRECTA PARA POS ---
AGENCIA_ORIGEN_POS_CORRECTA = "(8656) Retiro Corporativos"
SUBCUENTA_POS = "PEDIDOSYA SA -POS-"

def es_logistica_inversa_pos(row):
    """Detecta si un registro es logística inversa POS (origen incorrecto para esa subcuenta)"""
    subcuenta = str(row.get('Subcuenta', '')).strip()
    agencia_origen = str(row.get('Agencia origen', '')).strip()
    return subcuenta == SUBCUENTA_POS and agencia_origen != AGENCIA_ORIGEN_POS_CORRECTA

# --- INTERFAZ STREAMLIT ---
st.set_page_config(page_title="Calculadora de Lead Time", layout="wide")
st.title("📊 Calculadora de Lead Time - Indicadores Mejorados")
st.markdown("Sube tu reporte diario y obtén estadísticas + PPT listo para presentar.")
uploaded_file = st.file_uploader("📂 Sube tu archivo Excel", type=["xlsx", "xls"])

if uploaded_file is not None:
    try:
        df = pd.read_excel(uploaded_file, sheet_name="Prueba")
    except:
        df = pd.read_excel(uploaded_file, sheet_name=0)

    # --- Rellenar clientes vacíos con "EVENTUAL" ---
    if 'Cliente' in df.columns:
        df['Cliente'] = df['Cliente'].fillna("EVENTUAL")
        df['Cliente'] = df['Cliente'].apply(lambda x: "EVENTUAL" if str(x).strip() == "" else x)
        df['Cliente'] = df['Cliente'].astype(str)
    else:
        st.error("❌ La columna 'Cliente' no existe en el archivo.")
        st.stop()

    if 'Localidad destino' in df.columns:
        df['Loc'] = df['Localidad destino']

    # --- LIMPIAR LOCALIDADES ---
    if 'Loc' in df.columns:
        df['Loc'] = df['Loc'].apply(limpiar_localidad)

    # Convertir columnas de fecha
    df['Fecha'] = pd.to_datetime(df['Fecha'], errors='coerce')
    df['Fecha último estado'] = pd.to_datetime(df['Fecha último estado'], errors='coerce')

    df['Semana Calendario'] = df['Fecha'].apply(obtener_semana_calendario)

    # --- EXCLUSIONES CONSOLIDADAS ---
    exclusion_messages = []
    count_before_all = df.shape[0]

    if 'Agencia destino' in df.columns:
        count_before = df.shape[0]
        df = df[df['Agencia destino'] != "(6100) Administracion IPE"]
        excluded_count = count_before - df.shape[0]
        if excluded_count > 0:
            exclusion_messages.append(f"{excluded_count} guías con destino a '(6100) Administracion IPE' (paquetería interna)")

    if 'Loc' in df.columns:
        count_before = df.shape[0]
        df = df[~df['Loc'].str.upper().str.strip().eq("ADMINISTRACION BS AS, CAPITAL FEDERAL")]
        excluded_count = count_before - df.shape[0]
        if excluded_count > 0:
            exclusion_messages.append(f"{excluded_count} guías con destino a 'ADMINISTRACION BS AS, CAPITAL FEDERAL'")

    if 'Importe total' in df.columns:
        count_before = df.shape[0]
        df['Importe total'] = pd.to_numeric(df['Importe total'], errors='coerce')
        df = df[(df['Importe total'] != 0) & (df['Importe total'].notna())]
        excluded_count = count_before - df.shape[0]
        if excluded_count > 0:
            exclusion_messages.append(f"{excluded_count} guías con Importe total igual a 0 o nulo")

    if exclusion_messages:
        total_excluded = count_before_all - df.shape[0]
        exclusion_text = f"ℹ️ Se excluyeron {total_excluded} guías en total:\n"
        for i, msg in enumerate(exclusion_messages, 1):
            exclusion_text += f"  {i}. {msg}\n"
        exclusion_text += f"\n📊 Total restante para análisis: {df.shape[0]} guías"
        st.info(exclusion_text)

    # Determinar ZONA y Categoría
    df['ZONA'] = df['Loc'].apply(determinar_zona)
    df['Categoria'] = df['Loc'].apply(determinar_categoria)
    df.loc[df['Categoria'].str.startswith('AMBA', na=False), 'ZONA'] = 'AMBA'

    # --- DÍAS PROMETIDOS ---
    # Lista de provincias del interior para detectar origen
    PROVINCIAS_INTERIOR = [
        "SALTA", "TUCUMAN", "JUJUY", "CATAMARCA", "LA RIOJA",
        "SANTIAGO DEL ESTERO", "MENDOZA", "SAN JUAN", "SAN LUIS",
        "CORDOBA", "SANTA FE", "ENTRE RIOS", "CORRIENTES", "MISIONES",
        "CHACO", "FORMOSA", "NEUQUEN", "RIO NEGRO", "CHUBUT",
        "SANTA CRUZ", "TIERRA DEL FUEGO", "LA PAMPA"
    ]

    def determinar_dias_prometidos_robusta(row):
        """
        Determina los días prometidos según:
        - RIDERS: siempre 3 días.
        - POS con origen incorrecto: se marcará como logística inversa (no aplica SLA).
        - Origen interior → cualquier destino: 5 días.
        - Resto: según Categoría destino.
        """
        try:
            cliente = str(row.get('Cliente', '')).strip().upper()
            subcuenta = str(row.get('Subcuenta', '')).strip().upper()
            categoria = str(row.get('Categoria', '')).strip()
            origen = str(row.get('Origen', '')).strip().upper()

            # Excepción RIDERS
            if "DELIVERY HERO" in cliente and "RIDERS" in subcuenta:
                return 3

            # Extraer provincia del origen (formato "CIUDAD, PROVINCIA")
            provincia_origen = origen.split(",")[-1].strip() if "," in origen else origen
            origen_es_interior = any(prov in provincia_origen for prov in PROVINCIAS_INTERIOR)

            # Interior → cualquier destino: 5 días
            if origen_es_interior:
                return 5

            # AMBA/CABA origen → lógica por categoría destino
            if categoria == "AMBA cercano":
                return 2
            elif categoria == "AMBA interior":
                return 5
            elif categoria == "Buenos Aires interior":
                return 5
            else:
                return 5

        except Exception as e:
            return 5

    df['Días Prometidos'] = df.apply(determinar_dias_prometidos_robusta, axis=1)

    # --- CÁLCULO DE LEAD TIME ---
    def calcular_lead_time(row):
        try:
            estado = str(row['Estado']).lower()
            ed = str(row.get('ED', '')).upper() if 'ED' in df.columns else 'SI'
            entregado = (
                (ed == "NO" and "esperando retiro" in estado) or
                "entregada" in estado
            )
            fecha_actual_argentina = obtener_fecha_actual_argentina().replace(tzinfo=None)
            if entregado:
                lead_time = calcular_dias_habiles(row['Fecha'], row['Fecha último estado'])
            else:
                lead_time = calcular_dias_habiles(row['Fecha'], fecha_actual_argentina)
            return lead_time
        except Exception as e:
            return None

    df['Lead Time'] = df.apply(calcular_lead_time, axis=1)

    # --- CÁLCULO DE CUMPLIMIENTO ---
    def determinar_cumplimiento_mejorado(row):
        # --- EXCEPCIÓN LOGÍSTICA INVERSA POS ---
        # Si la subcuenta es POS pero el origen no es Retiro Corporativos,
        # es un retiro de logística inversa y no aplica al SLA de entregas
        if es_logistica_inversa_pos(row):
            return "Excluido - Logística Inversa POS"

        estado = str(row['Estado']).lower()
        ed = str(row.get('ED', '')).upper() if 'ED' in df.columns else 'SI'
        condicion_venta = str(row.get('Condición de venta', '')).upper() if 'Condición de venta' in df.columns else ''
        visitas = row.get('Visitas', 0) if 'Visitas' in df.columns else 0

        if "cancelada" in estado:
            return "Cancelada"
        
        if "contingencia" in estado:
            return "Contingencia"

        # --- CLIENTES EVENTUAL: detección por destinatario ---
        if row.get('Cliente', '') == "EVENTUAL":
            destinatario = ""
            if 'Destinatario' in row.index:
                destinatario_value = row['Destinatario']
                if pd.notna(destinatario_value) and str(destinatario_value).strip() != "":
                    destinatario = str(destinatario_value).lower().strip()
            palabras_devolucion = [
                "devolucion", "devolucion md", "devolucion p-ya", "dev. pedidos ya/",
                "devoluciones", "devo", "devol", "devolución", "devoluciónes",
                "devol pedido ya", "dev a origen"
            ]
            if destinatario and any(palabra in destinatario for palabra in palabras_devolucion):
                return "Devuelto"

        # --- DEVOLUCIONES POR ESTADO (APLICA A TODOS LOS CLIENTES) ---
        if ("devolución a remitente" in estado or "devuelta" in estado or
            "devolución informada" in estado or "devolucion informada" in estado or
            "devolución en destino" in estado):
            # Cumplido si hubo al menos una visita, sin importar los días
            if visitas > 0:
                return "Devuelto - Cumplido (Visita a Tiempo)"
            else:
                return "Devuelto"

        if ed == "NO" and "esperando retiro" in estado:
            if pd.notna(row['Lead Time']) and row['Lead Time'] <= row['Días Prometidos']:
                if condicion_venta == "PD":
                    return "Entregada - En Tiempo (PD: Pago Pendiente)"
                else:
                    return "Entregada - En Tiempo"
            else:
                if condicion_venta == "PD":
                    return "Entregada - Fuera de Tiempo (PD: Pago Pendiente)"
                else:
                    return "Entregada - Fuera de Tiempo"

        elif "entregada" in estado:
            if pd.notna(row['Lead Time']) and row['Lead Time'] <= row['Días Prometidos']:
                return "Entregada - En Tiempo"
            else:
                return "Entregada - Fuera de Tiempo"

        else:
            if pd.notna(row['Lead Time']):
                if row['Lead Time'] < row['Días Prometidos']:
                    base_estado = "Pendiente - En Tiempo"
                elif row['Lead Time'] == row['Días Prometidos']:
                    base_estado = "Pendiente - Último Día"
                else:
                    base_estado = "Pendiente - Fuera de Tiempo"

                estados_visita = [
                    "visita a domicilio", "reprogramada", "domicilio incompleto",
                    "domicilio incorrecto", "ausente", "rechazado"
                ]
                es_estado_visita = any(e in estado for e in estados_visita)

                if es_estado_visita and visitas > 0:
                    if "domicilio incompleto" in estado:
                        return base_estado + " (Datos Incompletos)"
                    elif "domicilio incorrecto" in estado:
                        return base_estado + " (Domicilio Incorrecto)"
                    elif "ausente" in estado:
                        return base_estado + " (Cliente Ausente)"
                    elif "rechazado" in estado:
                        return base_estado + " (Cliente Rechazó)"
                    else:
                        return base_estado + " (Visita Realizada)"
                else:
                    return base_estado
            else:
                return "Pendiente - Fuera de Tiempo"

    df['Cumplimiento'] = df.apply(determinar_cumplimiento_mejorado, axis=1)

    # --- CATEGORÍAS EXCLUIDAS DEL SLA ---
    EXCLUIDOS_SLA = ["Cancelada", "Excluido - Logística Inversa POS", "Devuelto - Cumplido (Visita a Tiempo)", "Contingencia"]

    def calcular_dias_restantes(row):
        cumplimiento = str(row['Cumplimiento'])
        if ("Pendiente" in cumplimiento and
            "Entregada" not in cumplimiento and
            "Devuelto" not in cumplimiento and
            "Cancelada" not in cumplimiento and
            "Excluido" not in cumplimiento):
            if pd.notna(row['Lead Time']):
                restantes = row['Días Prometidos'] - row['Lead Time']
                if restantes > 1:
                    return f"{int(restantes)} días restantes"
                elif restantes == 1:
                    return "Vence mañana"
                elif restantes == 0:
                    return "Vence hoy"
                else:
                    return "Vencido"
        return ""

    df['Días Restantes'] = df.apply(calcular_dias_restantes, axis=1)

    # --- ALERTAS ---

    def alerta_en_transito_demorado(row):
        try:
            estado = str(row['Estado']).lower()
            fecha_ultimo_estado = row['Fecha último estado']
            if "en tránsito" in estado and pd.notna(fecha_ultimo_estado):
                fecha_actual_argentina = obtener_fecha_actual_argentina().replace(tzinfo=None)
                dias_desde_ultimo_estado = calcular_dias_habiles(fecha_ultimo_estado, fecha_actual_argentina)
                if dias_desde_ultimo_estado is not None and dias_desde_ultimo_estado >= 5:
                    return "En tránsito demorado"
            return ""
        except Exception as e:
            return ""

    df['Alerta En Tránsito Demorado'] = df.apply(alerta_en_transito_demorado, axis=1)

    def alerta_creada_demorada(row):
        try:
            estado = str(row['Estado']).lower()
            fecha_ultimo_estado = row['Fecha último estado']
            fecha_creacion = row['Fecha']
            if "creada" in estado and pd.notna(fecha_ultimo_estado) and pd.notna(fecha_creacion):
                fecha_actual_argentina = obtener_fecha_actual_argentina().replace(tzinfo=None)
                diferencia_horas = (fecha_actual_argentina - fecha_ultimo_estado).total_seconds() / 3600
                if diferencia_horas >= 24:
                    dias_habiles_creada = calcular_dias_habiles(fecha_creacion, fecha_actual_argentina)
                    return f"Creada demorada ({diferencia_horas:.1f} horas, {dias_habiles_creada} días hábiles)"
                elif diferencia_horas >= 12:
                    return f"Creada próxima a vencer ({diferencia_horas:.1f} horas)"
            return ""
        except Exception as e:
            return ""

    df['Alerta Creada Demorada'] = df.apply(alerta_creada_demorada, axis=1)

    def alerta_seguimiento_visitas(row):
        try:
            estado = str(row['Estado']).lower()
            cumplimiento = str(row['Cumplimiento'])
            fecha_ultimo_estado = row['Fecha último estado']
            visitas = row.get('Visitas', 0) if 'Visitas' in df.columns else 0
            if visitas >= 2 and "Visita" in cumplimiento and pd.notna(fecha_ultimo_estado):
                fecha_actual_argentina = obtener_fecha_actual_argentina().replace(tzinfo=None)
                dias_desde_visita = calcular_dias_habiles(fecha_ultimo_estado, fecha_actual_argentina)
                if dias_desde_visita is not None and dias_desde_visita >= 3:
                    if "Datos Incompletos" in cumplimiento:
                        return "Solicitar datos completos"
                    elif "Domicilio Incorrecto" in cumplimiento:
                        return "Verificar domicilio"
                    elif "Cliente Ausente" in cumplimiento:
                        return "Coordinar nueva visita"
                    elif "Cliente Rechazó" in cumplimiento:
                        return "Sugerir devolución"
                    else:
                        return "Requiere seguimiento"
            return ""
        except Exception as e:
            return ""

    df['Alerta Seguimiento Visitas'] = df.apply(alerta_seguimiento_visitas, axis=1)

    def alerta_una_visita_sin_seguimiento(row):
        try:
            visitas = row.get('Visitas', 0) if 'Visitas' in df.columns else 0
            estado = str(row['Estado']).lower()
            fecha_ultimo_estado = row['Fecha último estado']
            cumplimiento = str(row['Cumplimiento'])
            estados_visita = [
                "visita a domicilio", "reprogramada", "domicilio incompleto",
                "domicilio incorrecto", "ausente", "rechazado"
            ]
            es_estado_visita = any(estado_visita in estado for estado_visita in estados_visita)
            if (visitas == 1 and
                es_estado_visita and
                pd.notna(fecha_ultimo_estado) and
                "Visita" in cumplimiento and
                "Devuelto" not in cumplimiento and
                "Entregada" not in cumplimiento):
                fecha_actual_argentina = obtener_fecha_actual_argentina().replace(tzinfo=None)
                dias_desde_visita = calcular_dias_habiles(fecha_ultimo_estado, fecha_actual_argentina)
                if dias_desde_visita is not None and dias_desde_visita >= 5:
                    return f"1 visita hace {dias_desde_visita} días hábiles - Sin seguimiento"
            return ""
        except Exception as e:
            return ""

    df['Alerta Una Visita Sin Seguimiento'] = df.apply(alerta_una_visita_sin_seguimiento, axis=1)

    def alerta_devolucion(row):
        try:
            estado = str(row['Estado']).lower()
            ed = str(row.get('ED', '')).upper() if 'ED' in df.columns else 'SI'
            fecha_ultimo_estado = row['Fecha último estado']
            if ed == "NO" and "esperando retiro" in estado and pd.notna(fecha_ultimo_estado):
                fecha_actual_argentina = obtener_fecha_actual_argentina().replace(tzinfo=None)
                dias_desde_ultimo_estado = calcular_dias_habiles(fecha_ultimo_estado, fecha_actual_argentina)
                if dias_desde_ultimo_estado is not None and dias_desde_ultimo_estado >= 15:
                    return "Sugerir devolución"
            return ""
        except Exception as e:
            return ""

    df['Alerta Devolución'] = df.apply(alerta_devolucion, axis=1)

    def alerta_redespacho(row):
        try:
            estado = str(row['Estado']).lower()
            fecha_ultimo_estado = row['Fecha último estado']
            if "redespachada" in estado and pd.notna(fecha_ultimo_estado):
                fecha_actual_argentina = obtener_fecha_actual_argentina().replace(tzinfo=None)
                dias_desde_ultimo_estado = calcular_dias_habiles(fecha_ultimo_estado, fecha_actual_argentina)
                if dias_desde_ultimo_estado is not None and dias_desde_ultimo_estado >= 2:
                    return "Redespacho demorado"
            return ""
        except Exception as e:
            return ""

    df['Alerta Redespacho'] = df.apply(alerta_redespacho, axis=1)

    def alerta_pendiente_fuera_tiempo(row):
        cumplimiento = str(row['Cumplimiento'])
        if cumplimiento == "Pendiente - Fuera de Tiempo":
            return "Fuera de tiempo crítico"
        return ""

    df['Alerta Pendiente Fuera Tiempo'] = df.apply(alerta_pendiente_fuera_tiempo, axis=1)

    def alerta_pago_pendiente(row):
        try:
            estado = str(row['Estado']).lower()
            condicion_venta = str(row.get('Condición de venta', '')).upper() if 'Condición de venta' in df.columns else ''
            fecha_ultimo_estado = row['Fecha último estado']
            if condicion_venta == "PD" and "esperando retiro" in estado and pd.notna(fecha_ultimo_estado):
                fecha_actual_argentina = obtener_fecha_actual_argentina().replace(tzinfo=None)
                dias_desde_ultimo_estado = calcular_dias_habiles(fecha_ultimo_estado, fecha_actual_argentina)
                if dias_desde_ultimo_estado is not None and dias_desde_ultimo_estado >= 5:
                    return "Pago pendiente demorado"
            return ""
        except Exception as e:
            return ""

    df['Alerta Pago Pendiente'] = df.apply(alerta_pago_pendiente, axis=1)

    def alerta_reprogramada_sin_visitas(row):
        try:
            estado = str(row['Estado']).lower()
            visitas = row.get('Visitas', 0) if 'Visitas' in df.columns else 0
            if "reprogramada" in estado and visitas == 0:
                fecha_actual_argentina = obtener_fecha_actual_argentina().replace(tzinfo=None)
                fecha_ultimo_estado = row['Fecha último estado']
                if pd.notna(fecha_ultimo_estado):
                    dias_desde_reprogramacion = calcular_dias_habiles(fecha_ultimo_estado, fecha_actual_argentina)
                    if dias_desde_reprogramacion is not None and dias_desde_reprogramacion >= 2:
                        return f"Reprogramada sin visita ({dias_desde_reprogramacion} días hábiles)"
                    else:
                        return "Reprogramada sin visita"
                else:
                    return "Reprogramada sin visita"
            return ""
        except Exception as e:
            return ""

    df['Alerta Reprogramada Sin Visitas'] = df.apply(alerta_reprogramada_sin_visitas, axis=1)

    def alerta_vencimiento_mañana(row):
        try:
            estado = str(row.get('Estado', '')).lower()
            cumplimiento = str(row.get('Cumplimiento', ''))
            lead_time = row.get('Lead Time')
            cliente = str(row.get('Cliente', '')).strip().upper()
            subcuenta = str(row.get('Subcuenta', '')).strip().upper()
            categoria = str(row.get('Categoria', '')).strip()
            fecha_creacion = row.get('Fecha')

            # No aplica a entregados, cancelados, devueltos ni logística inversa
            if ("entregada" in estado or
                "cancelada" in estado or
                "devuelto" in cumplimiento.lower() or
                "excluido" in cumplimiento.lower()):
                return ""

            if "DELIVERY HERO" in cliente and "RIDERS" in subcuenta:
                dias_prometidos_correcto = 3
            else:
                origen = str(row.get('Origen', '')).strip().upper()
                provincia_origen = origen.split(",")[-1].strip() if "," in origen else origen
                origen_es_interior = any(prov in provincia_origen for prov in PROVINCIAS_INTERIOR)
                if origen_es_interior:
                    dias_prometidos_correcto = 5
                elif categoria == "AMBA cercano":
                    dias_prometidos_correcto = 2
                elif categoria in ["AMBA interior", "Buenos Aires interior"]:
                    dias_prometidos_correcto = 5
                else:
                    dias_prometidos_correcto = 5

            if (pd.notna(lead_time) and isinstance(lead_time, (int, float)) and
                pd.notna(fecha_creacion)):
                if lead_time >= dias_prometidos_correcto:
                    return "Ya vencido"
                elif lead_time == dias_prometidos_correcto - 1:
                    fecha_actual = obtener_fecha_actual_argentina().date()
                    fecha_manana = fecha_actual + timedelta(days=1)
                    if es_dia_habil(fecha_manana):
                        return "Vence mañana"
                    else:
                        proximo_dia_habil = fecha_manana
                        while not es_dia_habil(proximo_dia_habil):
                            proximo_dia_habil += timedelta(days=1)
                        dias_hasta_proximo_habil = calcular_dias_habiles(fecha_actual, proximo_dia_habil)
                        if dias_hasta_proximo_habil == 1:
                            return f"Vence {proximo_dia_habil.strftime('%d/%m')}"
                        else:
                            return f"Vence en {dias_hasta_proximo_habil} días"
            return ""
        except Exception as e:
            return ""

    df['Alerta Vencimiento Mañana'] = df.apply(alerta_vencimiento_mañana, axis=1)

    # --- NUEVA ALERTA: ORIGEN RETIRO POS ---
    def alerta_origen_retiro_pos(row):
        try:
            if es_logistica_inversa_pos(row):
                return "Alerta Origen Retiro POS"
            return ""
        except:
            return ""

    df['Alerta Origen Retiro POS'] = df.apply(alerta_origen_retiro_pos, axis=1)

    # --- ASIGNAR PRIORIDAD ---
    def asignar_prioridad(row):
        if row['Alerta Vencimiento Mañana'] == "Ya vencido":
            return "ALTA - Ya Vencido"
        elif row['Alerta Pendiente Fuera Tiempo'] == "Fuera de tiempo crítico":
            return "ALTA - Fuera de Tiempo"
        elif row['Alerta Devolución'] == "Sugerir devolución":
            return "ALTA - Devolución Demorada"
        elif row['Alerta Redespacho'] == "Redespacho demorado":
            return "ALTA - Redespacho"
        elif row['Alerta En Tránsito Demorado'] != "":
            return "ALTA - Fuera de Tiempo"
        elif row['Alerta Reprogramada Sin Visitas'] != "":
            return "ALTA - Reprogramada Sin Visita"
        elif row['Alerta Creada Demorada'] != "" and "demorada" in row['Alerta Creada Demorada'].lower():
            return "ALTA - Creada Demorada"
        elif row['Alerta Vencimiento Mañana'] == "Vence mañana":
            return "ALTA - Vence Mañana"
        elif row['Alerta Seguimiento Visitas'] != "":
            return "MEDIA - Seguimiento Visitas"
        elif row['Alerta Una Visita Sin Seguimiento'] != "":
            return "MEDIA - 1 Visita Sin Seg."
        elif row['Alerta Creada Demorada'] != "" and "próxima a vencer" in row['Alerta Creada Demorada']:
            return "MEDIA - Creada Próxima a Vencer"
        elif row['Alerta Pago Pendiente'] == "Pago pendiente demorado":
            return "BAJA - Pago Pendiente"
        elif row['Alerta Origen Retiro POS'] != "":
            return "INFO - Logística Inversa POS"
        else:
            return ""

    df['Prioridad Alerta'] = df.apply(asignar_prioridad, axis=1)

    prioridad_orden = {
        "ALTA - Ya Vencido": 1,
        "ALTA - Fuera de Tiempo": 2,
        "ALTA - Devolución Demorada": 3,
        "ALTA - Redespacho": 4,
        "ALTA - Reprogramada Sin Visita": 5,
        "ALTA - Creada Demorada": 6,
        "ALTA - Vence Mañana": 7,
        "MEDIA - Seguimiento Visitas": 8,
        "MEDIA - 1 Visita Sin Seg.": 9,
        "MEDIA - Creada Próxima a Vencer": 10,
        "BAJA - Pago Pendiente": 11,
        "INFO - Logística Inversa POS": 12
    }
    df['Orden Prioridad'] = df['Prioridad Alerta'].map(prioridad_orden).fillna(999)
    df = df.sort_values('Orden Prioridad').reset_index(drop=True)

    # --- FILTROS ---
    st.sidebar.header("🔍 Filtros")

    if 'Cliente' in df.columns:
        clientes = sorted(df['Cliente'].dropna().unique())
        cliente_seleccionado = st.sidebar.selectbox("Cliente", ["Todos"] + clientes)
    else:
        st.error("❌ La columna 'Cliente' no existe en el archivo.")
        st.stop()

    df_filtrado = df.copy()
    if cliente_seleccionado != "Todos":
        df_filtrado = df_filtrado[df_filtrado['Cliente'] == cliente_seleccionado]

    if 'Subcuenta' in df_filtrado.columns:
        subcuentas = sorted(df_filtrado['Subcuenta'].dropna().unique())
        subcuenta_seleccionada = st.sidebar.selectbox("Subcuenta", ["Todas"] + subcuentas)
    else:
        st.error("❌ La columna 'Subcuenta' no existe en el archivo.")
        st.stop()

    if 'Agencia origen' in df_filtrado.columns:
        agencias_origen = sorted(df_filtrado['Agencia origen'].dropna().unique())
        agencia_origen_seleccionada = st.sidebar.selectbox("Agencia origen", ["Todas"] + agencias_origen)
    else:
        st.warning("⚠️ La columna 'Agencia origen' no existe. Se omitirá este filtro.")
        agencia_origen_seleccionada = "Todas"

    if 'Agencia destino' in df_filtrado.columns:
        agencias = sorted(df_filtrado['Agencia destino'].dropna().unique())
        agencia_seleccionada = st.sidebar.selectbox("Agencia destino", ["Todas"] + agencias)
    else:
        st.error("❌ La columna 'Agencia destino' no existe en el archivo.")
        st.stop()

    if 'ZONA' in df_filtrado.columns:
        zonas = sorted(df_filtrado['ZONA'].dropna().unique())
        zona_seleccionada = st.sidebar.selectbox("Zona", ["Todas"] + zonas)
    else:
        st.warning("⚠️ La columna 'ZONA' no existe. Se omitirá este filtro.")
        zona_seleccionada = "Todas"

    if 'Categoria' in df_filtrado.columns:
        todas_categorias = sorted(df_filtrado['Categoria'].dropna().unique())
        if zona_seleccionada == "AMBA":
            categorias_filtradas = [c for c in todas_categorias if c.startswith("AMBA")]
        elif zona_seleccionada == "INTERIOR":
            categorias_filtradas = [c for c in todas_categorias if not c.startswith("AMBA")]
        else:
            categorias_filtradas = todas_categorias
        categoria_seleccionada = st.sidebar.selectbox("Categoría", ["Todas"] + categorias_filtradas)
    else:
        st.warning("⚠️ La columna 'Categoria' no existe. Se omitirá este filtro.")
        categoria_seleccionada = "Todas"

    if 'ED' in df_filtrado.columns:
        ed_opciones = sorted(df_filtrado['ED'].dropna().unique())
        ed_seleccionada = st.sidebar.selectbox("Entrega a Domicilio (ED)", ["Todas"] + ed_opciones)
    else:
        st.warning("⚠️ La columna 'ED' no existe. Se omitirá este filtro.")
        ed_seleccionada = "Todas"

    if 'Condición de venta' in df_filtrado.columns:
        condiciones_venta = sorted(df_filtrado['Condición de venta'].dropna().unique())
        condicion_venta_seleccionada = st.sidebar.selectbox("Condición de venta", ["Todas"] + condiciones_venta)
    else:
        st.warning("⚠️ La columna 'Condición de venta' no existe. Se omitirá este filtro.")
        condicion_venta_seleccionada = "Todas"

    # Aplicar todos los filtros
    df_final = df.copy()
    if cliente_seleccionado != "Todos":
        df_final = df_final[df_final['Cliente'] == cliente_seleccionado]
    if subcuenta_seleccionada != "Todas":
        df_final = df_final[df_final['Subcuenta'] == subcuenta_seleccionada]
    if 'Agencia origen' in df_final.columns and agencia_origen_seleccionada != "Todas":
        df_final = df_final[df_final['Agencia origen'] == agencia_origen_seleccionada]
    if agencia_seleccionada != "Todas":
        df_final = df_final[df_final['Agencia destino'] == agencia_seleccionada]
    if 'ZONA' in df_final.columns and zona_seleccionada != "Todas":
        df_final = df_final[df_final['ZONA'] == zona_seleccionada]
    if 'Categoria' in df_final.columns and categoria_seleccionada != "Todas":
        df_final = df_final[df_final['Categoria'] == categoria_seleccionada]
    if 'ED' in df_final.columns and ed_seleccionada != "Todas":
        df_final = df_final[df_final['ED'] == ed_seleccionada]
    if 'Condición de venta' in df_final.columns and condicion_venta_seleccionada != "Todas":
        df_final = df_final[df_final['Condición de venta'] == condicion_venta_seleccionada]

    df = df_final

    if df.empty:
        st.warning("⚠️ No hay datos que coincidan con los filtros seleccionados.")
        st.stop()

    # --- PORCENTAJE DE CUMPLIMIENTO POR SEMANA ---
    st.header("📈 Porcentaje de Cumplimiento por Semana con Alertas de Variación")

    def calcular_cumplimiento_semana(grupo):
        total_pedidos_semana = grupo[~grupo['Cumplimiento'].isin(EXCLUIDOS_SLA)].shape[0]
        if total_pedidos_semana == 0:
            return 0
        cumplidos_semana = grupo[
            grupo['Cumplimiento'].isin([
                "Entregada - En Tiempo",
                "Entregada - En Tiempo (PD: Pago Pendiente)"
            ])
        ].shape[0]
        return (cumplidos_semana / total_pedidos_semana * 100)

    df_semana = df[~df['Cumplimiento'].isin(EXCLUIDOS_SLA)].groupby('Semana Calendario').apply(
        calcular_cumplimiento_semana
    ).reset_index(name='Porcentaje Cumplimiento')

    df_semana = df_semana.sort_values('Semana Calendario').reset_index(drop=True)
    df_semana['Variación vs Semana Anterior'] = df_semana['Porcentaje Cumplimiento'].diff()
    df_semana['Variación Porcentual'] = (df_semana['Porcentaje Cumplimiento'].pct_change() * 100).round(2)
    df_semana['Porcentaje Cumplimiento'] = df_semana['Porcentaje Cumplimiento'].round(2)

    def generar_alerta_variacion(row):
        variacion = row['Variación vs Semana Anterior']
        variacion_porcentual = row['Variación Porcentual']
        if pd.isna(variacion):
            return "🔵 Semana de referencia"
        elif variacion > 5:
            return f"🟢 Excelente! +{variacion:.1f}pts (+{variacion_porcentual:.1f}%)"
        elif variacion > 2:
            return f"🟡 Mejoró +{variacion:.1f}pts (+{variacion_porcentual:.1f}%)"
        elif variacion >= -2:
            return f"⚪ Estable {variacion:+.1f}pts"
        elif variacion > -5:
            return f"🟠 Alerta! {variacion:.1f}pts ({variacion_porcentual:+.1f}%)"
        else:
            return f"🔴 CRÍTICO! {variacion:.1f}pts ({variacion_porcentual:+.1f}%)"

    df_semana['Alerta Variación'] = df_semana.apply(generar_alerta_variacion, axis=1)

    st.subheader("Tabla de Cumplimiento por Semana con Alertas de Variación")
    st.dataframe(df_semana[['Semana Calendario', 'Porcentaje Cumplimiento', 'Alerta Variación']],
                use_container_width=True)

    if len(df_semana) > 1:
        fig_semana = px.line(
            df_semana, x='Semana Calendario', y='Porcentaje Cumplimiento',
            title='Evolución del Porcentaje de Cumplimiento por Semana',
            markers=True, line_shape='linear'
        )
        for i, row in df_semana.iterrows():
            if i > 0:
                variacion = row['Variación vs Semana Anterior']
                if abs(variacion) >= 2:
                    color = 'green' if variacion > 0 else 'red'
                    fig_semana.add_annotation(
                        x=row['Semana Calendario'], y=row['Porcentaje Cumplimiento'],
                        text=f"{variacion:+.1f}pts", showarrow=True, arrowhead=2,
                        arrowsize=1, arrowwidth=2, arrowcolor=color,
                        bgcolor=color, bordercolor=color,
                        font=dict(color='white', size=10)
                    )
        fig_semana.update_layout(
            xaxis_title='Semana Calendario', yaxis_title='Porcentaje de Cumplimiento (%)',
            yaxis=dict(range=[0, 100]), hovermode='x unified'
        )
        fig_semana.add_hline(y=80, line_dash="dash", line_color="red", annotation_text="Objetivo 80%")
        st.plotly_chart(fig_semana, use_container_width=True)

    st.subheader("📊 Resumen de Tendencias por Semana")
    if len(df_semana) > 1:
        ultima_semana = df_semana.iloc[-1]
        penultima_semana = df_semana.iloc[-2] if len(df_semana) > 1 else None
        mejora_semanas = (df_semana['Variación vs Semana Anterior'] > 0).sum()
        total_comparables = len(df_semana) - 1
        col1, col2, col3 = st.columns(3)
        with col1:
            if penultima_semana is not None:
                variacion_actual = ultima_semana['Variación vs Semana Anterior']
                if not pd.isna(variacion_actual):
                    if variacion_actual > 0:
                        st.success(f"📈 Semana {ultima_semana['Semana Calendario']}: **+{variacion_actual:.1f}pts** vs semana anterior")
                    else:
                        st.error(f"📉 Semana {ultima_semana['Semana Calendario']}: **{variacion_actual:.1f}pts** vs semana anterior")
        with col2:
            if total_comparables > 0:
                tasa_mejora = (mejora_semanas / total_comparables) * 100
                st.metric("📊 Tasa de Mejora Semanal", f"{tasa_mejora:.1f}%")
        with col3:
            if len(df_semana) >= 4:
                ultimas_4 = df_semana.tail(4)
                tendencia = ultimas_4['Porcentaje Cumplimiento'].mean()
                st.metric("📅 Promedio Últimas 4 Semanas", f"{tendencia:.1f}%")

    st.header("🔔 Notificaciones de Variación en Tiempo Real")
    if len(df_semana) > 1:
        ultima_semana = df_semana.iloc[-1]
        variacion_actual = ultima_semana['Variación vs Semana Anterior']
        if not pd.isna(variacion_actual):
            if variacion_actual < -10:
                st.error(f"🚨 **ALERTA CRÍTICA** — Caída drástica en la última semana: {variacion_actual:.1f} puntos. Revisar procesos urgentemente.")
            elif variacion_actual < -5:
                st.warning(f"⚠️ **ALERTA IMPORTANTE** — Caída significativa: {variacion_actual:.1f} puntos. Analizar causas.")
            elif variacion_actual > 10:
                st.success(f"🎉 **LOGRO DESTACADO** — Mejora excepcional: +{variacion_actual:.1f} puntos. Replicar buenas prácticas.")
            elif variacion_actual > 5:
                st.info(f"👍 **BUEN DESEMPEÑO** — Mejora significativa: +{variacion_actual:.1f} puntos. Mantener tendencia.")

    # --- CUMPLIMIENTO POR SEMANA Y ZONA ---
    st.header("📈 Porcentaje de Cumplimiento por Semana y Zona")

    def calcular_cumplimiento_semana_zona(grupo):
        total_pedidos_semana = grupo[~grupo['Cumplimiento'].isin(EXCLUIDOS_SLA)].shape[0]
        if total_pedidos_semana == 0:
            return 0
        cumplidos_semana = grupo[
            grupo['Cumplimiento'].isin([
                "Entregada - En Tiempo",
                "Entregada - En Tiempo (PD: Pago Pendiente)"
            ])
        ].shape[0]
        return (cumplidos_semana / total_pedidos_semana * 100)

    df_semana_zona = df[~df['Cumplimiento'].isin(EXCLUIDOS_SLA)].groupby(['Semana Calendario', 'ZONA']).apply(
        calcular_cumplimiento_semana_zona
    ).reset_index(name='Porcentaje Cumplimiento')

    df_semana_total = df[~df['Cumplimiento'].isin(EXCLUIDOS_SLA)].groupby('Semana Calendario').apply(
        calcular_cumplimiento_semana_zona
    ).reset_index(name='Porcentaje Cumplimiento')
    df_semana_total['ZONA'] = 'TOTAL'

    df_semana_completo = pd.concat([df_semana_zona, df_semana_total], ignore_index=True)
    df_semana_completo = df_semana_completo.sort_values(['Semana Calendario', 'ZONA']).reset_index(drop=True)
    df_semana_completo['Variación vs Semana Anterior'] = df_semana_completo.groupby('ZONA')['Porcentaje Cumplimiento'].diff()
    df_semana_completo['Variación Porcentual'] = (df_semana_completo.groupby('ZONA')['Porcentaje Cumplimiento'].pct_change() * 100).round(2)
    df_semana_completo['Porcentaje Cumplimiento'] = df_semana_completo['Porcentaje Cumplimiento'].round(2)

    def generar_alerta_variacion_zona(row):
        variacion = row['Variación vs Semana Anterior']
        variacion_porcentual = row['Variación Porcentual']
        if pd.isna(variacion):
            return "🔵 Semana de referencia"
        elif variacion > 5:
            return f"🟢 Excelente! +{variacion:.1f}pts (+{variacion_porcentual:.1f}%)"
        elif variacion > 2:
            return f"🟡 Mejoró +{variacion:.1f}pts (+{variacion_porcentual:.1f}%)"
        elif variacion >= -2:
            return f"⚪ Estable {variacion:+.1f}pts"
        elif variacion > -5:
            return f"🟠 Alerta! {variacion:.1f}pts ({variacion_porcentual:+.1f}%)"
        else:
            return f"🔴 CRÍTICO! {variacion:.1f}pts ({variacion_porcentual:+.1f}%)"

    df_semana_completo['Alerta Variación'] = df_semana_completo.apply(generar_alerta_variacion_zona, axis=1)

    st.subheader("Tabla de Cumplimiento por Semana y Zona con Alertas de Variación")
    df_pivot = df_semana_completo.pivot_table(
        index='Semana Calendario', columns='ZONA',
        values=['Porcentaje Cumplimiento', 'Alerta Variación'], aggfunc='first'
    )
    df_display = pd.DataFrame()
    for semana in df_pivot.index:
        row_data = {'Semana Calendario': semana}
        for zona in ['AMBA', 'INTERIOR', 'TOTAL']:
            if zona in df_pivot['Porcentaje Cumplimiento'].columns:
                row_data[f'{zona} - % Cumplimiento'] = df_pivot['Porcentaje Cumplimiento'][zona][semana]
                row_data[f'{zona} - Alerta'] = df_pivot['Alerta Variación'][zona][semana]
        df_display = pd.concat([df_display, pd.DataFrame([row_data])], ignore_index=True)
    df_display = df_display.sort_values('Semana Calendario').reset_index(drop=True)
    for col in df_display.columns:
        if '% Cumplimiento' in col:
            df_display[col] = df_display[col].apply(lambda x: f"{x:.1f}%" if pd.notna(x) else "N/A")
    st.dataframe(df_display, use_container_width=True)

    if len(df_semana_completo) > 1:
        fig_semana_zona = px.line(
            df_semana_completo, x='Semana Calendario', y='Porcentaje Cumplimiento',
            color='ZONA', title='Evolución del Porcentaje de Cumplimiento por Semana y Zona',
            markers=True, line_shape='linear',
            color_discrete_map={'AMBA': '#28a745', 'INTERIOR': '#007bff', 'TOTAL': '#ff6b00'}
        )
        for zona in ['AMBA', 'INTERIOR', 'TOTAL']:
            df_zona = df_semana_completo[df_semana_completo['ZONA'] == zona].sort_values('Semana Calendario')
            for i, row in df_zona.iterrows():
                if i > 0:
                    variacion = row['Variación vs Semana Anterior']
                    if not pd.isna(variacion) and abs(variacion) >= 2:
                        color = 'green' if variacion > 0 else 'red'
                        y_offset = 3 if zona == 'AMBA' else (-3 if zona == 'INTERIOR' else 0)
                        fig_semana_zona.add_annotation(
                            x=row['Semana Calendario'],
                            y=row['Porcentaje Cumplimiento'] + y_offset,
                            text=f"{variacion:+.1f}",
                            showarrow=True, arrowhead=2, arrowsize=1, arrowwidth=2,
                            arrowcolor=color, bgcolor=color, bordercolor=color,
                            font=dict(color='white', size=8),
                            yshift=10 if variacion > 0 else -10
                        )
        fig_semana_zona.update_layout(
            xaxis_title='Semana Calendario', yaxis_title='Porcentaje de Cumplimiento (%)',
            yaxis=dict(range=[0, 100]), hovermode='x unified',
            legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1)
        )
        fig_semana_zona.add_hline(y=80, line_dash="dash", line_color="red", annotation_text="Objetivo 80%")
        st.plotly_chart(fig_semana_zona, use_container_width=True)

    st.subheader("📊 Resumen de Tendencias por Semana y Zona")
    if len(df_semana_completo) > 1:
        zonas = ['AMBA', 'INTERIOR', 'TOTAL']
        cols = st.columns(3)
        for idx, zona in enumerate(zonas):
            with cols[idx]:
                st.subheader(f"Zona {zona}")
                df_zona = df_semana_completo[df_semana_completo['ZONA'] == zona].sort_values('Semana Calendario')
                if len(df_zona) > 1:
                    ultima_semana = df_zona.iloc[-1]
                    penultima_semana = df_zona.iloc[-2]
                    variacion_actual = ultima_semana['Variación vs Semana Anterior']
                    if not pd.isna(variacion_actual):
                        if variacion_actual > 0:
                            st.success(f"📈 **+{variacion_actual:.1f}pts** vs semana anterior")
                        else:
                            st.error(f"📉 **{variacion_actual:.1f}pts** vs semana anterior")
                    st.metric(f"Última Semana {ultima_semana['Semana Calendario']}", f"{ultima_semana['Porcentaje Cumplimiento']:.1f}%")
                    if len(df_zona) >= 4:
                        tendencia = df_zona.tail(4)['Porcentaje Cumplimiento'].mean()
                        st.metric("Promedio Últimas 4 Semanas", f"{tendencia:.1f}%")

    st.header("📊 Análisis Comparativo entre Zonas")
    if len(df_semana_completo) > 1:
        df_amba_sem = df_semana_completo[df_semana_completo['ZONA'] == 'AMBA'][['Semana Calendario', 'Porcentaje Cumplimiento']]
        df_interior_sem = df_semana_completo[df_semana_completo['ZONA'] == 'INTERIOR'][['Semana Calendario', 'Porcentaje Cumplimiento']]
        df_comparativo = pd.merge(df_amba_sem, df_interior_sem, on='Semana Calendario', suffixes=('_AMBA', '_INTERIOR'))
        df_comparativo['Diferencia (AMBA - INTERIOR)'] = df_comparativo['Porcentaje Cumplimiento_AMBA'] - df_comparativo['Porcentaje Cumplimiento_INTERIOR']
        fig_diferencias = px.bar(
            df_comparativo, x='Semana Calendario', y='Diferencia (AMBA - INTERIOR)',
            title='Diferencia de Cumplimiento: AMBA vs INTERIOR',
            color='Diferencia (AMBA - INTERIOR)',
            color_continuous_scale='RdYlGn', color_continuous_midpoint=0
        )
        fig_diferencias.update_layout(xaxis_title='Semana Calendario', yaxis_title='Diferencia (%)', hovermode='x unified')
        fig_diferencias.add_hline(y=0, line_dash="solid", line_color="black")
        st.plotly_chart(fig_diferencias, use_container_width=True)
        st.subheader("🔍 Resumen de Diferencias AMBA vs INTERIOR")
        if len(df_comparativo) > 0:
            ultima_diferencia = df_comparativo.iloc[-1]['Diferencia (AMBA - INTERIOR)']
            promedio_diferencia = df_comparativo['Diferencia (AMBA - INTERIOR)'].mean()
            col1, col2 = st.columns(2)
            with col1:
                if ultima_diferencia > 0:
                    st.success(f"**Última semana:** AMBA +{ultima_diferencia:.1f}pts sobre INTERIOR")
                else:
                    st.error(f"**Última semana:** INTERIOR +{abs(ultima_diferencia):.1f}pts sobre AMBA")
            with col2:
                if promedio_diferencia > 0:
                    st.info(f"**Promedio histórico:** AMBA +{promedio_diferencia:.1f}pts sobre INTERIOR")
                else:
                    st.warning(f"**Promedio histórico:** INTERIOR +{abs(promedio_diferencia):.1f}pts sobre AMBA")

    st.header("🔔 Notificaciones de Variación en Tiempo Real por Zona")
    if len(df_semana_completo) > 1:
        for zona in ['AMBA', 'INTERIOR', 'TOTAL']:
            df_zona = df_semana_completo[df_semana_completo['ZONA'] == zona].sort_values('Semana Calendario')
            if len(df_zona) > 1:
                ultima_semana = df_zona.iloc[-1]
                variacion_actual = ultima_semana['Variación vs Semana Anterior']
                if not pd.isna(variacion_actual):
                    if variacion_actual < -10:
                        st.error(f"🚨 **ALERTA CRÍTICA - {zona}** — Caída drástica: {variacion_actual:.1f} pts. Revisar urgentemente.")
                    elif variacion_actual < -5:
                        st.warning(f"⚠️ **ALERTA IMPORTANTE - {zona}** — Caída significativa: {variacion_actual:.1f} pts.")
                    elif variacion_actual > 10:
                        st.success(f"🎉 **LOGRO DESTACADO - {zona}** — Mejora excepcional: +{variacion_actual:.1f} pts.")
                    elif variacion_actual > 5:
                        st.info(f"👍 **BUEN DESEMPEÑO - {zona}** — Mejora significativa: +{variacion_actual:.1f} pts.")

    # Agregar columnas de semana al df principal
    mapeo_semana = df_semana.set_index('Semana Calendario')['Porcentaje Cumplimiento'].to_dict()
    mapeo_alerta = df_semana.set_index('Semana Calendario')['Alerta Variación'].to_dict()
    mapeo_variacion = df_semana.set_index('Semana Calendario')['Variación vs Semana Anterior'].to_dict()
    df['Porcentaje Cumplimiento Semana'] = df['Semana Calendario'].map(mapeo_semana)
    df['Alerta Variación Semana'] = df['Semana Calendario'].map(mapeo_alerta)
    df['Variación vs Semana Anterior'] = df['Semana Calendario'].map(mapeo_variacion)

    columnas = df.columns.tolist()
    pos_semana = columnas.index('Semana Calendario')
    nuevas_columnas = ['Porcentaje Cumplimiento Semana', 'Alerta Variación Semana', 'Variación vs Semana Anterior']
    for i, col in enumerate(nuevas_columnas):
        columnas.insert(pos_semana + 1 + i, col)
        columnas.remove(col)
    df = df[columnas]

    df['Porcentaje Cumplimiento Semana'] = df['Porcentaje Cumplimiento Semana'].apply(
        lambda x: f"{x:.1f}%" if pd.notna(x) else "N/A"
    )
    df['Variación vs Semana Anterior'] = df['Variación vs Semana Anterior'].apply(
        lambda x: f"{x:+.1f} pts" if pd.notna(x) else "N/A"
    )

    # --- ALERTAS CRÍTICAS SIDEBAR ---
    st.sidebar.header("🚨 Alertas Críticas de Variación")
    if len(df_semana) > 1:
        alertas_criticas = df_semana[
            (df_semana['Variación vs Semana Anterior'] < -5) &
            (pd.notna(df_semana['Variación vs Semana Anterior']))
        ]
        if not alertas_criticas.empty:
            st.sidebar.error("### 📉 Caídas Significativas")
            for _, alerta in alertas_criticas.iterrows():
                st.sidebar.write(f"**Semana {alerta['Semana Calendario']}**: {alerta['Variación vs Semana Anterior']:.1f}pts")
        mejoras_significativas = df_semana[
            (df_semana['Variación vs Semana Anterior'] > 5) &
            (pd.notna(df_semana['Variación vs Semana Anterior']))
        ]
        if not mejoras_significativas.empty:
            st.sidebar.success("### 📈 Mejoras Significativas")
            for _, mejora in mejoras_significativas.iterrows():
                st.sidebar.write(f"**Semana {mejora['Semana Calendario']}**: +{mejora['Variación vs Semana Anterior']:.1f}pts")

    # --- ESTADÍSTICAS ---
    st.header("📊 Indicadores/Alertas")

    total_pedidos = df[~df['Cumplimiento'].isin(EXCLUIDOS_SLA)].shape[0]
    entregados = df[df['Cumplimiento'].str.startswith("Entregada")].shape[0]
    devueltos = df[df['Cumplimiento'] == "Devuelto"].shape[0]                     # Solo los que NO cumplieron
    devuelto_cumplido_visita = df[df['Cumplimiento'] == "Devuelto - Cumplido (Visita a Tiempo)"].shape[0]
    total_devueltos = devueltos + devuelto_cumplido_visita                        # <-- Total real de devueltos
    canceladas = df[df['Cumplimiento'] == "Cancelada"].shape[0]
    logistica_inversa_pos_count = df[df['Cumplimiento'] == "Excluido - Logística Inversa POS"].shape[0]
    pendientes_reales = total_pedidos - entregados - devueltos   # devueltos aquí son los que sí cuentan para SLA

    visita_en_tiempo = df[df['Cumplimiento'].str.contains("Visita en Tiempo", na=False)].shape[0]
    visita_fuera_tiempo = df[df['Cumplimiento'] == "Pendiente - Visita Fuera de Tiempo"].shape[0]

    en_tiempo = df[df['Cumplimiento'] == "Entregada - En Tiempo"].shape[0]
    en_tiempo_pd = df[df['Cumplimiento'] == "Entregada - En Tiempo (PD: Pago Pendiente)"].shape[0]
    fuera_tiempo = df[df['Cumplimiento'] == "Entregada - Fuera de Tiempo"].shape[0]
    fuera_tiempo_pd = df[df['Cumplimiento'] == "Entregada - Fuera de Tiempo (PD: Pago Pendiente)"].shape[0]
    devuelto_count = total_devueltos
    devuelto_sin_cumplir = devuelto_count - devuelto_cumplido_visita
    pendiente_en_tiempo = df[df['Cumplimiento'] == "Pendiente - En Tiempo"].shape[0]
    pendiente_fuera_tiempo = df[df['Cumplimiento'] == "Pendiente - Fuera de Tiempo"].shape[0]
    pendiente_ultimo_dia = df[df['Cumplimiento'] == "Pendiente - Último Día"].shape[0]

    sla_principal = ((en_tiempo + en_tiempo_pd) / total_pedidos * 100) if total_pedidos > 0 else 0
    cumplimiento_tradicional = ((en_tiempo + en_tiempo_pd) / entregados * 100) if entregados > 0 else 0
    cumplimiento_gestion = ((en_tiempo + en_tiempo_pd + visita_en_tiempo) / total_pedidos * 100) if total_pedidos > 0 else 0

    total_visitas = visita_en_tiempo + visita_fuera_tiempo
    efectividad_visitas = (visita_en_tiempo / total_visitas * 100) if total_visitas > 0 else 0

    primer_intento_entrega = df[
        (df['Cumplimiento'].str.startswith("Entregada")) &
        (df.get('Visitas', 0) <= 1)
    ].shape[0]
    fadr = (primer_intento_entrega / entregados * 100) if entregados > 0 else 0

    total_visitas_entregados = df[df['Cumplimiento'].str.startswith("Entregada")]['Visitas'].sum()
    pedidos_con_visita = df[(df['Cumplimiento'].str.startswith("Entregada")) & (df.get('Visitas', 0) >= 1)].shape[0]
    pedidos_por_visita = (pedidos_con_visita / total_visitas_entregados) if total_visitas_entregados > 0 else 0

    def es_rechazo_ausente_regex(estado):
        if pd.isna(estado):
            return False
        estado_str = str(estado).strip()
        patrones = [
            r'\[Motivo:\s*(Rechazado|Ausente)',
            r'\[Motivo:\s*.*(rechaz|ausent)',
            r'(rechazado|ausente).*\[Motivo:',
            r'cliente\s+(rechazó|no aceptó|ausente|no se presentó)',
            r'motivo.*rechaz|motivo.*ausent'
        ]
        for patron in patrones:
            if re.search(patron, estado_str, re.IGNORECASE):
                return True
        return False

    rechazos_ausentes = df[
        df['Estado'].apply(es_rechazo_ausente_regex) &
        (df['Visitas'] > 0)
    ].shape[0]
    total_con_visita = df[df['Visitas'] > 0].shape[0]
    tasa_rechazo_ausencia = (rechazos_ausentes / total_con_visita * 100) if total_con_visita > 0 else 0

    alertas_creada_criticas = df[df['Alerta Creada Demorada'].str.contains("demorada", na=False)].shape[0]
    alertas_creada_preventivas = df[df['Alerta Creada Demorada'].str.contains("próxima a vencer", na=False)].shape[0]

    col1, col2, col3, col4, col5 = st.columns(5)
    with col1:
        st.metric("📦 Total Pedidos (Excl. Canceladas/Inversa)", total_pedidos)
        st.metric("🎯 SLA Principal", f"{sla_principal:.1f}%")
    with col2:
        st.metric("✅ Entregados", entregados, f"{(entregados/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%")
        st.metric("📊 Cumplimiento Entregas", f"{cumplimiento_tradicional:.1f}%")
    with col3:
        st.metric("🔄 Devueltos", devueltos, f"{(devueltos/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%")
        st.metric("📋 Cumplimiento Gestión", f"{cumplimiento_gestion:.1f}%")
    with col4:
        st.metric("⏳ Pendientes", pendientes_reales, f"{(pendientes_reales/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%")
        st.metric("🚫 Tasa Rechazo/Ausencia", f"{tasa_rechazo_ausencia:.1f}%")
    with col5:
        st.metric("🏭 Logística Inversa POS", logistica_inversa_pos_count)
        st.metric("🚨 Creadas Demoradas (>24h)", alertas_creada_criticas)

    if st.sidebar.checkbox("🔍 Mostrar pedidos con rechazo/ausencia detectados"):
        ejemplos_rechazo = df[df['Estado'].apply(es_rechazo_ausente_regex) & (df['Visitas'] > 0)]
        if not ejemplos_rechazo.empty:
            st.sidebar.write(f"📋 Ejemplos detectados ({len(ejemplos_rechazo)}):")
            st.sidebar.dataframe(ejemplos_rechazo[['Guia', 'Estado', 'Visitas']].head(5))
        else:
            st.sidebar.info("No se detectaron pedidos con rechazo/ausencia")

    # --- TABLA DE RESUMEN ---
    st.header("📈 Detalle de Estados")
    resumen_data = {
        "Categoría": [
            "TOTAL PEDIDOS (Excl. Canceladas/Logística Inversa POS)",
            "ENTREGADOS",
            " - En Tiempo",
            " - En Tiempo (PD)",
            " - Fuera de Tiempo",
            " - Fuera de Tiempo (PD)",
            "DEVUELTOS",                    # Solo los que penalizan
            "CANCELADAS",
            "EXCLUIDOS - Logística Inversa POS",
            "PENDIENTES CON VISITA",
            " - Visita en Tiempo",
            " - Visita Fuera de Tiempo",
            "PENDIENTES SIN VISITA",
            " - En Tiempo",
            " - Último Día",
            " - Fuera de Tiempo",
            "SLA PRINCIPAL (En Tiempo/Total)",
            "TASA RECHAZO/AUSENCIA"
        ],
        "Cantidad": [
            total_pedidos, entregados, en_tiempo, en_tiempo_pd,
            fuera_tiempo, fuera_tiempo_pd,
            devueltos,                      # ← Solo los que penalizan
            canceladas,
            logistica_inversa_pos_count,
            visita_en_tiempo + visita_fuera_tiempo,
            visita_en_tiempo, visita_fuera_tiempo,
            pendiente_en_tiempo + pendiente_ultimo_dia + pendiente_fuera_tiempo,
            pendiente_en_tiempo, pendiente_ultimo_dia, pendiente_fuera_tiempo,
            "", rechazos_ausentes
        ],
        "Porcentaje": [
            "100%",
            f"{(entregados/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%",
            f"{(en_tiempo/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%",
            f"{(en_tiempo_pd/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%",
            f"{(fuera_tiempo/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%",
            f"{(fuera_tiempo_pd/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%",
            f"{(devueltos/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%",   # ← sobre base gestionable
            f"{(canceladas/(total_pedidos + canceladas + logistica_inversa_pos_count)*100):.1f}%" if (total_pedidos + canceladas + logistica_inversa_pos_count) > 0 else "0%",
            f"{(logistica_inversa_pos_count/(total_pedidos + canceladas + logistica_inversa_pos_count)*100):.1f}%" if (total_pedidos + canceladas + logistica_inversa_pos_count) > 0 else "0%",
            f"{((visita_en_tiempo + visita_fuera_tiempo)/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%",
            f"{(visita_en_tiempo/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%",
            f"{(visita_fuera_tiempo/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%",
            f"{((pendiente_en_tiempo + pendiente_ultimo_dia + pendiente_fuera_tiempo)/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%",
            f"{(pendiente_en_tiempo/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%",
            f"{(pendiente_ultimo_dia/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%",
            f"{(pendiente_fuera_tiempo/total_pedidos*100):.1f}%" if total_pedidos > 0 else "0%",
            f"{sla_principal:.1f}%",
            f"{tasa_rechazo_ausencia:.1f}%"
        ]
    }
    resumen_df = pd.DataFrame(resumen_data)
    st.dataframe(resumen_df, use_container_width=True)

    # --- DESGLOSE AMBA ---
    st.header("📍 Desglose AMBA: Cercano vs Interior")
    df_amba = df[(df['Categoria'].str.startswith('AMBA', na=False)) & (~df['Cumplimiento'].isin(EXCLUIDOS_SLA))]
    if not df_amba.empty:
        df_cercano = df_amba[df_amba['Categoria'] == 'AMBA cercano']
        total_cercano = len(df_cercano)
        en_tiempo_cercano = df_cercano[df_cercano['Cumplimiento'].isin([
            "Entregada - En Tiempo", "Entregada - En Tiempo (PD: Pago Pendiente)"
        ])].shape[0]
        pct_cercano = (en_tiempo_cercano / total_cercano * 100) if total_cercano > 0 else 0

        df_interior_amba = df_amba[df_amba['Categoria'] == 'AMBA interior']
        total_interior_amba = len(df_interior_amba)
        en_tiempo_interior_amba = df_interior_amba[df_interior_amba['Cumplimiento'].isin([
            "Entregada - En Tiempo", "Entregada - En Tiempo (PD: Pago Pendiente)"
        ])].shape[0]
        pct_interior_amba = (en_tiempo_interior_amba / total_interior_amba * 100) if total_interior_amba > 0 else 0

        data_amba = {
            "Categoría": ["AMBA Cercano", "AMBA Interior"],
            "Total Gestionables": [total_cercano, total_interior_amba],
            "Entregas a Tiempo": [en_tiempo_cercano, en_tiempo_interior_amba],
            "% Cumplimiento": [f"{pct_cercano:.1f}%", f"{pct_interior_amba:.1f}%"],
            "Días Prometidos": ["2 días", "5 días"]
        }
        df_amba_display = pd.DataFrame(data_amba)
        st.dataframe(df_amba_display, use_container_width=True)
        fig_amba = px.bar(
            df_amba_display, x="Categoría", y="Total Gestionables",
            text="Total Gestionables", color="Categoría",
            title="Volumen de envíos AMBA por subcategoría",
            color_discrete_map={"AMBA Cercano": "#28a745", "AMBA Interior": "#fd7e14"}
        )
        fig_amba.update_traces(texttemplate='%{text}', textposition='outside')
        st.plotly_chart(fig_amba, use_container_width=True)
        col1, col2 = st.columns(2)
        with col1:
            st.metric("🎯 Cumplimiento AMBA Cercano", f"{pct_cercano:.1f}%",
                     delta=f"{en_tiempo_cercano} de {total_cercano} pedidos")
        with col2:
            st.metric("🎯 Cumplimiento AMBA Interior", f"{pct_interior_amba:.1f}%",
                     delta=f"{en_tiempo_interior_amba} de {total_interior_amba} pedidos")
    else:
        st.info("No hay datos de AMBA para mostrar con los filtros actuales.")

    # --- SECCIÓN: CONTEO DIARIO AMBA ---
    st.header("📆 Conteo Diario AMBA (Cercano / Interior)")
    st.markdown("Desglose diario de entregas a tiempo, fuera de tiempo y pendientes.")

    try:
        conteo_amba, detalle_pendientes = calcular_conteo_diario_amba(df)

        if not conteo_amba.empty:
            # Mostrar tabla resumen
            st.subheader("📊 Resumen diario")
            # Formatear porcentajes solo para visualización
            conteo_amba_display = conteo_amba.copy()
            conteo_amba_display['% Entregas en Tiempo'] = conteo_amba_display['% Entregas en Tiempo'].apply(lambda x: f"{x:.1%}")
            conteo_amba_display['% Entregas Totales'] = conteo_amba_display['% Entregas Totales'].apply(lambda x: f"{x:.1%}")
            st.dataframe(conteo_amba_display, use_container_width=True)

            # Detalle de pendientes por día
            if not detalle_pendientes.empty:
                st.subheader("📋 Detalle de pendientes por día")
                for fecha, grupo in detalle_pendientes.groupby('Fecha'):
                    with st.expander(f"Pendientes del {fecha} (creados el día hábil anterior)"):
                        st.dataframe(grupo.drop(columns=['Fecha']), use_container_width=True)

            # Generar Excel con formato de porcentaje
            output_conteo = io.BytesIO()
            with pd.ExcelWriter(output_conteo, engine='openpyxl') as writer:
                # Hoja resumen
                conteo_amba.to_excel(writer, sheet_name='Resumen Diario', index=False)
                # Formatear columnas de porcentaje en el archivo
                ws_resumen = writer.sheets['Resumen Diario']
                for col in ['H', 'I']:  # Columnas % Entregas en Tiempo y % Entregas Totales
                    for cell in ws_resumen[col]:
                        if isinstance(cell.value, (int, float)):
                            cell.number_format = '0.0%'

                # Hoja detalle
                if not detalle_pendientes.empty:
                    detalle_pendientes.to_excel(writer, sheet_name='Detalle Pendientes', index=False)

            output_conteo.seek(0)

            st.download_button(
                label="📥 Descargar Conteo Diario AMBA (Excel)",
                data=output_conteo,
                file_name="Conteo_Diario_AMBA.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )
        else:
            st.info("No hay datos de AMBA que coincidan con los filtros actuales.")
    except Exception as e:
        st.error(f"Error al generar el conteo diario AMBA: {e}") 

    # --- GRÁFICO DE TORTA (CON SLICE LOGÍSTICA INVERSA POS) ---
    categorias_mejoradas = [
        "Entregada - En Tiempo",
        "Entregada - En Tiempo (PD)",
        "Entregada - Fuera de Tiempo",
        "Entregada - Fuera de Tiempo (PD)",
        "Devuelto - Cumplido (Visita a Tiempo)",
        "Devuelto (sin visita / fuera de plazo)",
        "Cancelada",
        "Excluido - Logística Inversa POS",
        "Pendiente - Visita en Tiempo",
        "Pendiente - Visita Fuera de Tiempo",
        "Pendiente - En Tiempo",
        "Pendiente - Último Día",
        "Pendiente - Fuera de Tiempo"
    ]
    valores_mejorados = [
        en_tiempo, en_tiempo_pd, fuera_tiempo, fuera_tiempo_pd,
        devuelto_cumplido_visita, devuelto_sin_cumplir, # ← en este orden
        canceladas, logistica_inversa_pos_count,
        visita_en_tiempo, visita_fuera_tiempo,
        pendiente_en_tiempo, pendiente_ultimo_dia, pendiente_fuera_tiempo
    ]
    fig1 = px.pie(
        names=categorias_mejoradas,
        values=valores_mejorados,
        title="Distribución de Cumplimiento (Entregas vs Retiros vs Canceladas)",
        color=categorias_mejoradas,
        color_discrete_map={
            "Entregada - En Tiempo": "#28a745",
            "Entregada - En Tiempo (PD)": "#2ecc71",
            "Entregada - Fuera de Tiempo": "#dc3545",
            "Entregada - Fuera de Tiempo (PD)": "#e74c3c",
            "Devuelto - Cumplido (Visita a Tiempo)": "#8e44ad",
            "Devuelto (sin visita / fuera de plazo)": "#9b59b6",
            "Cancelada": "#95a5a6",
            "Excluido - Logística Inversa POS": "#17a2b8",
            "Pendiente - Visita en Tiempo": "#3498db",
            "Pendiente - Visita Fuera de Tiempo": "#e67e22",
            "Pendiente - En Tiempo": "#ffc107",
            "Pendiente - Último Día": "#fd7e14",
            "Pendiente - Fuera de Tiempo": "#6c757d"
        },
        hole=0.4
    )
    fig1.update_traces(textinfo='percent+value', textposition='inside')
    st.plotly_chart(fig1, use_container_width=True)

    # --- COMPARATIVA DE INDICADORES ---
    st.header("📈 Comparativa de Indicadores de Cumplimiento")
    indicadores_comparativa = {
        "Indicador": ["SLA Principal", "Cumplimiento Entregas", "Cumplimiento Gestión", "Tasa Rechazo/Ausencia"],
        "Porcentaje": [sla_principal, cumplimiento_tradicional, cumplimiento_gestion, tasa_rechazo_ausencia],
        "Descripción": [
            "Entregas en tiempo / Total pedidos",
            "Entregas en tiempo / Total entregados",
            "Gestión total (entregas + visitas en tiempo)",
            "Rechazos y ausencias / Total con visita"
        ]
    }
    df_comparativa = pd.DataFrame(indicadores_comparativa)
    fig2 = px.bar(
        df_comparativa, x="Indicador", y="Porcentaje", color="Indicador",
        text="Porcentaje", title="Comparativa de Diferentes Indicadores de Cumplimiento",
        color_discrete_map={
            "SLA Principal": "#28a745", "Cumplimiento Entregas": "#ffc107",
            "Cumplimiento Gestión": "#007bff", "Tasa Rechazo/Ausencia": "#dc3545"
        }
    )
    fig2.update_traces(texttemplate='%{y:.1f}%', textposition='outside')
    fig2.update_layout(showlegend=False)
    st.plotly_chart(fig2, use_container_width=True)

    # --- CUMPLIMIENTO POR CLIENTE ---
    st.header("📈 Cumplimiento Real por Cliente")

    def calcular_entregas_en_tiempo(grupo):
        return grupo[grupo['Cumplimiento'].isin(["Entregada - En Tiempo", "Entregada - En Tiempo (PD: Pago Pendiente)"])].shape[0]

    def calcular_visitas_en_tiempo(grupo):
        return grupo[grupo['Cumplimiento'].str.contains("Visita en Tiempo", na=False)].shape[0]

    df_cliente = df[~df['Cumplimiento'].isin(EXCLUIDOS_SLA)].groupby('Cliente').agg(
        Total_Pedidos=('Guia', 'count'),
        Entregas_En_Tiempo=('Cumplimiento', lambda x: calcular_entregas_en_tiempo(x.to_frame().assign(Cumplimiento=x))),
        Visitas_En_Tiempo=('Cumplimiento', lambda x: calcular_visitas_en_tiempo(x.to_frame().assign(Cumplimiento=x)))
    ).reset_index()
    df_cliente['Total_Pedidos'] = pd.to_numeric(df_cliente['Total_Pedidos'], errors='coerce').fillna(0)
    df_cliente['Entregas_En_Tiempo'] = pd.to_numeric(df_cliente['Entregas_En_Tiempo'], errors='coerce').fillna(0)
    df_cliente['Visitas_En_Tiempo'] = pd.to_numeric(df_cliente['Visitas_En_Tiempo'], errors='coerce').fillna(0)
    df_cliente['Cumplimiento_Real'] = ((df_cliente['Entregas_En_Tiempo'] + df_cliente['Visitas_En_Tiempo']) / df_cliente['Total_Pedidos'].replace(0, 1) * 100).round(2)
    df_cliente['Cumplimiento_Real'] = df_cliente['Cumplimiento_Real'].replace([np.inf, -np.inf], 0).fillna(0)
    df_cliente = df_cliente[df_cliente['Total_Pedidos'] >= 5]
    fig_cliente = px.bar(
        df_cliente.sort_values('Cumplimiento_Real', ascending=True),
        x='Cumplimiento_Real', y='Cliente', orientation='h',
        text='Cumplimiento_Real', title='Cumplimiento Real por Cliente (Mín. 5 pedidos)',
        color='Cumplimiento_Real', color_continuous_scale='RdYlGn'
    )
    fig_cliente.update_traces(texttemplate='%{text:.1f}%', textposition='outside')
    fig_cliente.update_layout(yaxis={'categoryorder': 'total ascending'})
    st.plotly_chart(fig_cliente, use_container_width=True)

    # --- CUMPLIMIENTO POR ZONA ---
    st.header("🗺️ Cumplimiento Real por Zona (AMBA vs INTERIOR)")
    df_zona = df[~df['Cumplimiento'].isin(EXCLUIDOS_SLA)].groupby('ZONA').agg(
        Total_Pedidos=('Guia', 'count'),
        Entregas_En_Tiempo=('Cumplimiento', lambda x: calcular_entregas_en_tiempo(x.to_frame().assign(Cumplimiento=x))),
        Visitas_En_Tiempo=('Cumplimiento', lambda x: calcular_visitas_en_tiempo(x.to_frame().assign(Cumplimiento=x)))
    ).reset_index()
    df_zona['Total_Pedidos'] = pd.to_numeric(df_zona['Total_Pedidos'], errors='coerce').fillna(0)
    df_zona['Entregas_En_Tiempo'] = pd.to_numeric(df_zona['Entregas_En_Tiempo'], errors='coerce').fillna(0)
    df_zona['Visitas_En_Tiempo'] = pd.to_numeric(df_zona['Visitas_En_Tiempo'], errors='coerce').fillna(0)
    df_zona['Cumplimiento_Real'] = ((df_zona['Entregas_En_Tiempo'] + df_zona['Visitas_En_Tiempo']) / df_zona['Total_Pedidos'].replace(0, 1) * 100).round(2)
    df_zona['Cumplimiento_Real'] = df_zona['Cumplimiento_Real'].replace([np.inf, -np.inf], 0).fillna(0)
    fig_zona = px.bar(
        df_zona, x='ZONA', y='Cumplimiento_Real', text='Cumplimiento_Real',
        title='Comparativa de Cumplimiento Real por Zona',
        color='ZONA', color_discrete_map={'AMBA': '#28a745', 'INTERIOR': '#007bff'}
    )
    fig_zona.update_traces(texttemplate='%{y:.1f}%', textposition='outside')
    st.plotly_chart(fig_zona, use_container_width=True)

    # --- TOP 5 AGENCIAS CON MÁS CANCELACIONES ---
    if canceladas > 0:
        st.header("📉 Top 5 Agencias Origen con Más Cancelaciones")
        top_agencias_cancel = df[df['Cumplimiento'] == "Cancelada"]['Agencia origen'].value_counts().head(5)
        if not top_agencias_cancel.empty:
            fig_cancel = px.bar(
                top_agencias_cancel, x=top_agencias_cancel.values, y=top_agencias_cancel.index,
                orientation='h', text=top_agencias_cancel.values,
                title="Top 5 Agencias Origen con Más Cancelaciones",
                color=top_agencias_cancel.values, color_continuous_scale='Blues'
            )
            fig_cancel.update_traces(texttemplate='%{text}', textposition='outside')
            fig_cancel.update_layout(yaxis={'categoryorder': 'total ascending'})
            st.plotly_chart(fig_cancel, use_container_width=True)
    else:
        st.info("✅ No hay cancelaciones para mostrar.")

    # --- TOP 5 LOCALIDADES FUERA DE TIEMPO ---
    st.header("⏳ Top 5 Localidades (Loc) con Más Pedidos Fuera de Tiempo")
    top_loc_fuera_tiempo = df[df['Cumplimiento'] == "Pendiente - Fuera de Tiempo"]['Loc'].value_counts().head(5)
    if not top_loc_fuera_tiempo.empty:
        fig_loc = px.bar(
            top_loc_fuera_tiempo, x=top_loc_fuera_tiempo.values, y=top_loc_fuera_tiempo.index,
            orientation='h', text=top_loc_fuera_tiempo.values,
            title="Top 5 Localidades con Más Pedidos Fuera de Tiempo",
            color=top_loc_fuera_tiempo.values, color_continuous_scale='Reds'
        )
        fig_loc.update_traces(texttemplate='%{text}', textposition='outside')
        fig_loc.update_layout(yaxis={'categoryorder': 'total ascending'})
        st.plotly_chart(fig_loc, use_container_width=True)
    else:
        st.info("✅ No hay pedidos 'Fuera de Tiempo' para mostrar.")

    # =============================================
    # SECCIONES DE ALERTAS
    # =============================================

    # --- ALERTA: LOGÍSTICA INVERSA POS ---
    alertas_origen_pos = df[df['Alerta Origen Retiro POS'] != ""]
    if not alertas_origen_pos.empty:
        st.header("🏭 Alertas de Origen Retiro POS")
        st.write(f"Pedidos de **{SUBCUENTA_POS}** con origen distinto a **{AGENCIA_ORIGEN_POS_CORRECTA}** (logística inversa excluida del SLA):")
        columnas_alerta = [
            'Guia', 'Importe total', 'Cliente', 'Subcuenta',
            'Agencia origen', 'Origen', 'Destinatario', 'Loc',
            'ZONA', 'Fecha', 'Estado', 'Cumplimiento',
            'Alerta Origen Retiro POS'
        ]
        columnas_existentes = [col for col in columnas_alerta if col in df.columns]
        df_alerta = alertas_origen_pos[columnas_existentes]
        st.dataframe(df_alerta)
        excel_data = generar_excel_desde_df(df_alerta, "Alertas Origen Retiro POS")
        st.download_button(
            label="📥 Descargar Alertas Origen Retiro POS (Excel)",
            data=excel_data,
            file_name="Alertas_Origen_Retiro_POS.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # --- ALERTA: ESTADO CREADA DEMORADO ---
    alertas_creada_demorada_df = df[df['Alerta Creada Demorada'] != ""]
    if not alertas_creada_demorada_df.empty:
        st.header("🚨 Alertas de Estado 'Creada' Demorado")
        st.write("Los siguientes pedidos están en estado 'Creada' por más de 24 horas:")
        columnas_alerta = [
            'Guia', 'Importe total', 'Cliente', 'Destinatario', 'Loc', 'ZONA',
            'Fecha', 'Fecha último estado', 'Estado',
            'Alerta Creada Demorada', 'Prioridad Alerta'
        ]
        columnas_existentes = [col for col in columnas_alerta if col in df.columns]
        df_alerta = alertas_creada_demorada_df[columnas_existentes]
        st.dataframe(df_alerta)
        excel_data = generar_excel_desde_df(df_alerta, "Alertas Creada Demorada")
        st.download_button(
            label="📥 Descargar Alertas de Estado 'Creada' Demorado (Excel)",
            data=excel_data,
            file_name="Alertas_Creada_Demorada.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # --- ALERTA: SEGUIMIENTO VISITAS (2+) ---
    alertas_seguimiento = df[df['Alerta Seguimiento Visitas'] != ""]
    if not alertas_seguimiento.empty:
        st.header("🔄 Alertas de Seguimiento de Visitas (2+ Visitas)")
        st.write("Pedidos con múltiples visitas que requieren acción:")
        columnas_alerta = ['Guia', 'Importe total', 'Cliente', 'Destinatario', 'Tel Destinatario',
                          'Loc', 'ZONA', 'Estado', 'Visitas', 'Fecha último estado',
                          'Cumplimiento', 'Alerta Seguimiento Visitas', 'Prioridad Alerta']
        columnas_existentes = [col for col in columnas_alerta if col in df.columns]
        df_alerta = alertas_seguimiento[columnas_existentes]
        st.dataframe(df_alerta)
        excel_data = generar_excel_desde_df(df_alerta, "Alertas Seguimiento Visitas")
        st.download_button(
            label="📥 Descargar Alertas de Seguimiento de Visitas (Excel)",
            data=excel_data,
            file_name="Alertas_Seguimiento_Visitas.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # --- ALERTA: UNA VISITA SIN SEGUIMIENTO ---
    alertas_una_visita = df[df['Alerta Una Visita Sin Seguimiento'] != ""]
    if not alertas_una_visita.empty:
        st.header("⏰ Alertas de Una Visita Sin Seguimiento")
        st.write("Pedidos con solo una visita sin seguimiento en 5+ días hábiles:")
        columnas_alerta = ['Guia', 'Importe total', 'Cliente', 'Destinatario', 'Tel Destinatario',
                          'Loc', 'ZONA', 'Estado', 'Visitas', 'Fecha último estado',
                          'Cumplimiento', 'Alerta Una Visita Sin Seguimiento', 'Prioridad Alerta']
        columnas_existentes = [col for col in columnas_alerta if col in df.columns]
        df_alerta = alertas_una_visita[columnas_existentes]
        st.dataframe(df_alerta)
        excel_data = generar_excel_desde_df(df_alerta, "Alertas Una Visita Sin Seguimiento")
        st.download_button(
            label="📥 Descargar Alertas Una Visita Sin Seguimiento (Excel)",
            data=excel_data,
            file_name="Alertas_Una_Visita_Sin_Seguimiento.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # --- ALERTA: DEVOLUCIÓN ---
    alertas_devolucion = df[df['Alerta Devolución'] == "Sugerir devolución"]
    if not alertas_devolucion.empty:
        st.header("🚨 Alertas de Devolución")
        st.write("Pedidos en 'Esperando retiro' por más de 15 días hábiles:")
        columnas_alerta = ['Guia', 'Importe total', 'Cliente', 'Destinatario', 'Tel Destinatario',
                          'Loc', 'ZONA', 'Fecha último estado', 'Estado', 'Alerta Devolución', 'Prioridad Alerta']
        columnas_existentes = [col for col in columnas_alerta if col in df.columns]
        df_alerta = alertas_devolucion[columnas_existentes]
        st.dataframe(df_alerta)
        excel_data = generar_excel_desde_df(df_alerta, "Alertas Devolución")
        st.download_button(
            label="📥 Descargar Alertas Devolución (Excel)",
            data=excel_data,
            file_name="Alertas_Devolucion.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # --- ALERTA: REDESPACHO ---
    alertas_redespacho = df[df['Alerta Redespacho'] == "Redespacho demorado"]
    if not alertas_redespacho.empty:
        st.header("🚨 Alertas de Redespacho Demorado")
        st.write("Pedidos en estado 'Redespacho' por más de 48 horas hábiles:")
        columnas_alerta = ['Guia', 'Importe total', 'Cliente', 'Destinatario', 'Tel Destinatario',
                          'Loc', 'ZONA', 'Fecha último estado', 'Estado', 'Alerta Redespacho', 'Prioridad Alerta']
        columnas_existentes = [col for col in columnas_alerta if col in df.columns]
        df_alerta = alertas_redespacho[columnas_existentes]
        st.dataframe(df_alerta)
        excel_data = generar_excel_desde_df(df_alerta, "Alertas Redespacho")
        st.download_button(
            label="📥 Descargar Alertas Redespacho (Excel)",
            data=excel_data,
            file_name="Alertas_Redespacho.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # --- ALERTA: REPROGRAMADA SIN VISITAS ---
    alertas_reprogramada_sin_visitas = df[df['Alerta Reprogramada Sin Visitas'] != ""]
    if not alertas_reprogramada_sin_visitas.empty:
        st.header("🚨 Alertas de Reprogramada Sin Visitas")
        st.write("Pedidos en estado 'Reprogramada' sin visitas registradas:")
        columnas_alerta = ['Guia', 'Importe total', 'Cliente', 'Destinatario', 'Tel Destinatario',
                          'Loc', 'ZONA', 'Estado', 'Visitas', 'Fecha último estado',
                          'Cumplimiento', 'Alerta Reprogramada Sin Visitas', 'Prioridad Alerta']
        columnas_existentes = [col for col in columnas_alerta if col in df.columns]
        df_alerta = alertas_reprogramada_sin_visitas[columnas_existentes]
        st.dataframe(df_alerta)
        excel_data = generar_excel_desde_df(df_alerta, "Alertas Reprogramada Sin Visitas")
        st.download_button(
            label="📥 Descargar Alertas Reprogramada Sin Visitas (Excel)",
            data=excel_data,
            file_name="Alertas_Reprogramada_Sin_Visitas.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # --- ALERTA: EN TRÁNSITO DEMORADO ---
    alertas_en_transito = df[df['Alerta En Tránsito Demorado'] != ""]
    if not alertas_en_transito.empty:
        st.header("🚨 Alertas de En Tránsito Demorado")
        st.write("Pedidos en estado 'En tránsito' por más de 48 horas hábiles:")
        columnas_alerta = ['Guia', 'Importe total', 'Cliente', 'Destinatario', 'Tel Destinatario',
                          'Loc', 'ZONA', 'Fecha último estado', 'Estado',
                          'Alerta En Tránsito Demorado', 'Prioridad Alerta']
        columnas_existentes = [col for col in columnas_alerta if col in df.columns]
        df_alerta = alertas_en_transito[columnas_existentes]
        st.dataframe(df_alerta)
        excel_data = generar_excel_desde_df(df_alerta, "Alertas En Tránsito Demorado")
        st.download_button(
            label="📥 Descargar Alertas En Tránsito Demorado (Excel)",
            data=excel_data,
            file_name="Alertas_En_Transito_Demorado.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # --- ALERTA: PENDIENTE FUERA DE TIEMPO ---
    alertas_pendiente_fuera_tiempo = df[df['Alerta Pendiente Fuera Tiempo'] == "Fuera de tiempo crítico"]
    if not alertas_pendiente_fuera_tiempo.empty:
        st.header("🚨 Alertas de Pendiente Fuera de Tiempo")
        st.write("Pedidos pendientes fuera del tiempo de entrega prometido:")
        columnas_alerta = ['Guia', 'Importe total', 'Cliente', 'Destinatario', 'Tel Destinatario',
                          'Loc', 'ZONA', 'Fecha último estado', 'Estado', 'Días Prometidos',
                          'Lead Time', 'Alerta Pendiente Fuera Tiempo', 'Prioridad Alerta']
        columnas_existentes = [col for col in columnas_alerta if col in df.columns]
        df_alerta = alertas_pendiente_fuera_tiempo[columnas_existentes]
        st.dataframe(df_alerta)
        excel_data = generar_excel_desde_df(df_alerta, "Alertas Pendiente Fuera Tiempo")
        st.download_button(
            label="📥 Descargar Alertas Pendiente Fuera Tiempo (Excel)",
            data=excel_data,
            file_name="Alertas_Pendiente_Fuera_Tiempo.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # --- ALERTA: PAGO PENDIENTE ---
    alertas_pago_pendiente = df[df['Alerta Pago Pendiente'] == "Pago pendiente demorado"]
    if not alertas_pago_pendiente.empty:
        st.header("🚨 Alertas de Pago Pendiente Demorado")
        st.write("Pedidos con condición PD en 'Esperando retiro' por más de 5 días hábiles:")
        columnas_alerta = ['Guia', 'Importe total', 'Cliente', 'Destinatario', 'Tel Destinatario',
                          'Loc', 'ZONA', 'Fecha último estado', 'Estado', 'Condición de venta',
                          'Alerta Pago Pendiente', 'Prioridad Alerta']
        columnas_existentes = [col for col in columnas_alerta if col in df.columns]
        df_alerta = alertas_pago_pendiente[columnas_existentes]
        st.dataframe(df_alerta)
        excel_data = generar_excel_desde_df(df_alerta, "Alertas Pago Pendiente")
        st.download_button(
            label="📥 Descargar Alertas Pago Pendiente (Excel)",
            data=excel_data,
            file_name="Alertas_Pago_Pendiente.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # --- ALERTA: VENCIMIENTO ---
    alertas_vencimiento_mañana_df = df[df['Alerta Vencimiento Mañana'].isin(["Vence mañana", "Ya vencido"])]
    if not alertas_vencimiento_mañana_df.empty:
        st.header("🚨 Alertas de Vencimiento")
        st.write("Pedidos que **vencen mañana** o que **ya están vencidos**:")
        vence_mañana_count = len(alertas_vencimiento_mañana_df[alertas_vencimiento_mañana_df['Alerta Vencimiento Mañana'] == "Vence mañana"])
        ya_vencido_count = len(alertas_vencimiento_mañana_df[alertas_vencimiento_mañana_df['Alerta Vencimiento Mañana'] == "Ya vencido"])
        col1, col2 = st.columns(2)
        with col1:
            st.metric("📅 Vencen Mañana", vence_mañana_count)
        with col2:
            st.metric("⏰ Ya Vencidos", ya_vencido_count)
        columnas_alerta = [
            'Guia', 'Importe total', 'Cliente', 'Subcuenta', 'Destinatario', 'Tel Destinatario',
            'Loc', 'ZONA', 'Fecha', 'Fecha último estado', 'Estado',
            'Días Prometidos', 'Lead Time', 'Cumplimiento', 'Días Restantes',
            'Alerta Vencimiento Mañana', 'Prioridad Alerta'
        ]
        columnas_existentes = [col for col in columnas_alerta if col in alertas_vencimiento_mañana_df.columns]
        df_alerta = alertas_vencimiento_mañana_df[columnas_existentes]
        st.dataframe(df_alerta)
        excel_data = generar_excel_desde_df(df_alerta, "Alertas Vencimiento")
        st.download_button(
            label="📥 Descargar Alertas de Vencimiento (Excel)",
            data=excel_data,
            file_name="Alertas_Vencimiento.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # --- DESCARGA COMBINADA ---
    st.header("📥 Descarga Combinada de Todas las Alertas")
    todas_alertas = df[
        (df['Alerta Seguimiento Visitas'] != "") |
        (df['Alerta Una Visita Sin Seguimiento'] != "") |
        (df['Alerta Devolución'] == "Sugerir devolución") |
        (df['Alerta Redespacho'] == "Redespacho demorado") |
        (df['Alerta Pendiente Fuera Tiempo'] == "Fuera de tiempo crítico") |
        (df['Alerta Pago Pendiente'] == "Pago pendiente demorado") |
        (df['Alerta En Tránsito Demorado'] != "") |
        (df['Alerta Creada Demorada'] != "") |
        (df['Alerta Vencimiento Mañana'].isin(["Vence mañana", "Ya vencido"])) |
        (df['Alerta Origen Retiro POS'] != "")
    ]
    if not todas_alertas.empty:
        columnas_todas = [
            'Guia', 'Importe total', 'Cliente', 'Agencia origen', 'Subcuenta',
            'Destinatario', 'Tel Destinatario', 'Loc', 'Agencia destino', 'ZONA',
            'Visitas', 'Fecha último estado', 'Días Prometidos', 'Lead Time',
            'Estado', 'Cumplimiento', 'Prioridad Alerta',
            'Alerta Seguimiento Visitas', 'Alerta Una Visita Sin Seguimiento',
            'Alerta Devolución', 'Alerta Redespacho',
            'Alerta Pendiente Fuera Tiempo', 'Alerta Pago Pendiente',
            'Alerta Vencimiento Mañana', 'Alerta En Tránsito Demorado',
            'Alerta Creada Demorada', 'Alerta Origen Retiro POS'
        ]
        columnas_existentes = [col for col in columnas_todas if col in df.columns]
        df_todas = todas_alertas[columnas_existentes]
        st.dataframe(df_todas)
        excel_todas = generar_excel_desde_df(df_todas, "Todas las Alertas")
        st.download_button(
            label="📥 Descargar Todas las Alertas (Excel)",
            data=excel_todas,
            file_name="Todas_Alertas_Combinadas.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
    else:
        st.info("✅ No hay alertas activas en este momento.")

    # --- VISTA PREVIA ---
    st.header("🔍 Vista Previa de Datos con Alertas de Variación")
    columnas_mostrar = [
        'Cliente', 'Subcuenta', 'Agencia origen', 'Agencia destino', 'Condición de venta',
        'Fecha', 'Semana Calendario', 'Porcentaje Cumplimiento Semana',
        'Alerta Variación Semana', 'Variación vs Semana Anterior',
        'Fecha último estado', 'Estado', 'Visitas', 'ED', 'Loc', 'ZONA', 'Categoría', 'Producto',
        'Lead Time', 'Días Prometidos', 'Cumplimiento', 'Días Restantes', 'Prioridad Alerta',
        'Alerta Seguimiento Visitas', 'Alerta Una Visita Sin Seguimiento',
        'Alerta Devolución', 'Alerta Redespacho', 'Alerta Pendiente Fuera Tiempo',
        'Alerta Pago Pendiente', 'Alerta En Tránsito Demorado', 'Alerta Creada Demorada',
        'Alerta Vencimiento Mañana', 'Alerta Origen Retiro POS'
    ]
    columnas_existentes = [col for col in columnas_mostrar if col in df.columns]
    df_vista_previa = df[columnas_existentes].head(10)
    st.dataframe(df_vista_previa)

    # --- DESCARGAS GENERALES ---
    st.header("📥 Descargas Generales Actualizadas")

    if 'Categoria' in df.columns:
        df.rename(columns={'Categoria': 'Categoría'}, inplace=True)

    alertas_vencimiento_count = len(df[df['Alerta Vencimiento Mañana'].isin(["Vence mañana", "Ya vencido"])])
    vence_mañana_count = len(df[df['Alerta Vencimiento Mañana'] == "Vence mañana"])
    ya_vencido_count = len(df[df['Alerta Vencimiento Mañana'] == "Ya vencido"])

    stats_data = {
        "Métrica": [
            "Total Pedidos (Excl. Canceladas/Logística Inversa POS)",
            "Entregados", "Devueltos", "Canceladas",
            "Logística Inversa POS (Excluidos SLA)", "Pendientes Reales",
            "Entregada - En Tiempo", "Entregada - En Tiempo (PD)",
            "Entregada - Fuera de Tiempo", "Entregada - Fuera de Tiempo (PD)",
            "Devuelto", "Cancelada", "Logística Inversa POS",
            "Pendiente - Visita en Tiempo", "Pendiente - Visita Fuera de Tiempo",
            "Pendiente - En Tiempo", "Pendiente - Último Día", "Pendiente - Fuera de Tiempo",
            "SLA Principal (%)", "Cumplimiento Entregas (%)", "Cumplimiento Gestión (%)",
            "FADR (%)", "Pedidos por Visita", "Tasa Rechazo/Ausencia (%)",
            "Alertas Creada Demoradas", "Alertas Creada Próximas a Vencer",
            "Alertas Vencimiento Total", "Alertas Vencen Mañana", "Alertas Ya Vencidos",
            "Alertas Logística Inversa POS"
        ],
        "Valor": [
            total_pedidos, entregados, devueltos, canceladas,
            logistica_inversa_pos_count, pendientes_reales,
            en_tiempo, en_tiempo_pd, fuera_tiempo, fuera_tiempo_pd,
            devuelto_count, canceladas, logistica_inversa_pos_count,
            visita_en_tiempo, visita_fuera_tiempo,
            pendiente_en_tiempo, pendiente_ultimo_dia, pendiente_fuera_tiempo,
            f"{sla_principal:.2f}%", f"{cumplimiento_tradicional:.2f}%", f"{cumplimiento_gestion:.2f}%",
            f"{fadr:.2f}%", f"{pedidos_por_visita:.2f}", f"{tasa_rechazo_ausencia:.2f}%",
            alertas_creada_criticas, alertas_creada_preventivas,
            alertas_vencimiento_count, vence_mañana_count, ya_vencido_count,
            len(alertas_origen_pos)
        ]
    }

    stats_df = pd.DataFrame(stats_data)

    output_excel_actualizado = io.BytesIO()
    with pd.ExcelWriter(output_excel_actualizado, engine='openpyxl') as writer:
        df.to_excel(writer, sheet_name="Base Completa", index=False)
        df_semana.to_excel(writer, sheet_name="Cumplimiento por Semana", index=False)
        stats_df.to_excel(writer, sheet_name="Estadísticas", index=False)
    output_excel_actualizado.seek(0)

    col_btn1, col_btn2 = st.columns(2)
    with col_btn1:
        st.download_button(
            label="📥 Descargar Excel Actualizado (Completo)",
            data=output_excel_actualizado,
            file_name="Reporte_LeadTime_Con_Porcentaje_Semana.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
    excel_vista_actualizada = generar_excel_desde_df(df[columnas_existentes], "Vista Previa Completa")
    with col_btn2:
        st.download_button(
            label="📥 Descargar Vista Previa Actualizada",
            data=excel_vista_actualizada,
            file_name="Vista_Previa_Actualizada.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # --- POWERPOINT ---
    def crear_pptx():
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Reporte de Cumplimiento de Entregas"
        subtitle.text = "Lead Time - Indicadores Mejorados\nGenerado automáticamente"

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Resumen Ejecutivo"
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Métricas Clave:"
        p.font.bold = True
        p.font.size = Pt(20)

        alertas_vencimiento_count_ppt = len(df[df['Alerta Vencimiento Mañana'].isin(["Vence mañana", "Ya vencido"])])
        vence_mañana_count_ppt = len(df[df['Alerta Vencimiento Mañana'] == "Vence mañana"])
        ya_vencido_count_ppt = len(df[df['Alerta Vencimiento Mañana'] == "Ya vencido"])

        metrics = [
            f"• Total de pedidos (Excl. Canceladas/Inversa): {total_pedidos}",
            f"• Entregados: {entregados} ({(entregados/total_pedidos*100):.1f}%)" if total_pedidos > 0 else f"• Entregados: {entregados}",
            f"• Devueltos: {devueltos} ({(devueltos/total_pedidos*100):.1f}%)" if total_pedidos > 0 else f"• Devueltos: {devueltos}",
            f"• Canceladas: {canceladas}",
            f"• Logística Inversa POS (excluidos SLA): {logistica_inversa_pos_count}",
            f"• SLA Principal: {sla_principal:.1f}%",
            f"• Cumplimiento Entregas: {cumplimiento_tradicional:.1f}%",
            f"• Cumplimiento Gestión: {cumplimiento_gestion:.1f}%",
            f"• FADR (1er Intento): {fadr:.1f}%",
            f"• Tasa Rechazo/Ausencia: {tasa_rechazo_ausencia:.1f}%",
            f"• Alertas Activas: {len(todas_alertas)}",
            f"• Alertas Creada Demoradas: {alertas_creada_criticas}",
            f"• Alertas Vencimiento: {alertas_vencimiento_count_ppt}",
            f"  - Vencen mañana: {vence_mañana_count_ppt}",
            f"  - Ya vencidos: {ya_vencido_count_ppt}",
            f"• Alertas Logística Inversa POS: {len(alertas_origen_pos)}"
        ]
        for metric in metrics:
            p = tf.add_paragraph()
            p.text = metric
            p.font.size = Pt(16)

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Cumplimiento por Semana"
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Evolución Semanal:"
        p.font.bold = True
        p.font.size = Pt(18)
        if len(df_semana) > 0:
            ultimas_semanas = df_semana.tail(4)
            for _, semana in ultimas_semanas.iterrows():
                p = tf.add_paragraph()
                p.text = f"Semana {semana['Semana Calendario']}: {semana['Porcentaje Cumplimiento']:.1f}% - {semana['Alerta Variación']}"
                p.font.size = Pt(14)

        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        return pptx_buffer

    if st.button("📊 Generar y Descargar PowerPoint Actualizado"):
        pptx_data = crear_pptx()
        st.download_button(
            label="⬇️ Descargar Presentación PPTX Actualizada",
            data=pptx_data,
            file_name="Reporte_LeadTime_Presentacion_Actualizada.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

else:
    st.info("👆 Por favor, sube un archivo Excel para comenzar.")
    st.markdown("""
    **Instrucciones:**
    1. Haz clic en "Browse files".
    2. Selecciona tu archivo Excel.
    3. ¡Listo! La app calculará automáticamente y mostrará gráficos y botones de descarga.
    """)